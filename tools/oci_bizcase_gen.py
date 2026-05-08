#!/usr/bin/env python3
"""
OCI Deal Accelerator — Business Case Deck Generator (.pptx)

Produces an 8-10 slide business case deck using the Oracle FY26 official
PowerPoint template layouts. Target audience: customer CxO / decision-makers.

Usage:
    python oci_bizcase_gen.py --spec business-case.yaml --output business-case.pptx

Or import and use programmatically:
    from oci_bizcase_gen import BusinessCaseDeckGenerator
    gen = BusinessCaseDeckGenerator.from_spec(spec)
    gen.save("business-case.pptx")
"""

import yaml
import argparse
import re
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

try:
    from oci_pptx_base import Colors, Layouts, OraclePresBase
    from oci_business_case_model import (
        enrich_adbs_to_adbd_business_case,
        write_bom_outputs,
    )
    from pptx_validator import validate_pptx
except ModuleNotFoundError:
    from tools.oci_pptx_base import Colors, Layouts, OraclePresBase
    from tools.oci_business_case_model import (
        enrich_adbs_to_adbd_business_case,
        write_bom_outputs,
    )
    from tools.pptx_validator import validate_pptx


# ============================================================
# Defensive input helpers
# ============================================================

def pick(mapping: dict, *keys, default=""):
    """Return the first non-empty value for any of the provided keys."""
    if not isinstance(mapping, dict):
        return default
    for key in keys:
        if key in mapping:
            value = mapping.get(key)
            if value is not None:
                return value
    return default


def pick_list(mapping: dict, *keys):
    """Return the first list-like field normalized to a list."""
    if not isinstance(mapping, dict):
        return []
    for key in keys:
        if key in mapping:
            value = mapping.get(key)
            if value is None:
                return []
            if isinstance(value, list):
                return value
            return [value]
    return []


def _timeline_to_weeks(text: str) -> int | None:
    """Parse a simple duration phrase like '12 weeks' or '6 months'."""
    if not text:
        return None
    value = str(text).lower()
    match = re.search(r"(\d+)\s*(week|weeks|month|months)", value)
    if not match:
        return None
    amount = int(match.group(1))
    unit = match.group(2)
    return amount * 4 if unit.startswith("month") else amount


def _extract_business_case_text(bc: dict) -> str:
    parts = [
        pick(bc, "executive_summary", "summary"),
        pick(bc, "primary_driver", "main_driver"),
    ]
    drivers = bc.get("drivers", {})
    if isinstance(drivers, dict):
        parts.extend([
            pick(drivers, "primary", "primary_driver", "main_driver"),
            pick(drivers, "urgency", "why_now"),
        ])
        coi = drivers.get("cost_of_inaction", {}) or {}
        if isinstance(coi, dict):
            parts.extend(coi.values())
    return " ".join(str(p) for p in parts if p)


def _infer_driver_headline(text: str) -> str:
    lowered = text.lower()
    if any(token in lowered for token in ("budget", "cost", "tco", "license", "refresh")):
        return "Reduce cost and avoid infrastructure refresh"
    if any(token in lowered for token in ("dr", "disaster recovery", "resilien", "rto", "rpo")):
        return "Improve resilience and business continuity"
    if any(token in lowered for token in ("pci", "regulator", "compliance", "data residency", "sovereignty")):
        return "Address compliance and sovereignty requirements"
    return "Modernize the platform on OCI"


def _infer_driver_objective(text: str) -> str:
    """Return a noun phrase suitable for recommendation prose."""
    lowered = text.lower()
    if any(token in lowered for token in ("budget", "cost", "tco", "license", "refresh")):
        return "the cost-reduction and refresh-avoidance case"
    if any(token in lowered for token in ("dr", "disaster recovery", "resilien", "rto", "rpo")):
        return "the resilience and business-continuity posture"
    if any(token in lowered for token in ("pci", "regulator", "compliance", "data residency", "sovereignty")):
        return "the compliance and sovereignty position"
    return "the OCI modernization approach"


def _infer_business_case_drivers(text: str) -> dict:
    headline = _infer_driver_headline(text)
    urgency = "Use the current planning window to validate scope, commercials, and delivery readiness."
    if _timeline_to_weeks(text):
        urgency = f"Align the case with the stated delivery window ({_timeline_to_weeks(text)} weeks)."
    return {
        "primary": headline,
        "urgency": urgency,
        "cost_of_inaction": {
            "financial": "Continuing with the current platform preserves avoidable run costs and refresh exposure.",
            "operational": "The delivery team keeps spending time on manual operations instead of modernization work.",
            "strategic": "Delaying the decision extends technical debt and slows business change.",
        },
    }


def _infer_business_case_risks(text: str) -> dict:
    lowered = text.lower()
    migration = [
        {
            "risk": "Migration scope may be larger than the initial discovery captured.",
            "mitigation": "Baseline applications, data stores, and dependencies before locking the plan.",
        },
        {
            "risk": "Operating model changes may require team enablement before go-live.",
            "mitigation": "Define the target RACI and training plan during design.",
        },
    ]
    if any(token in lowered for token in ("dr", "disaster recovery", "rto", "rpo")):
        migration.append({
            "risk": "Resilience objectives may require rehearsal and explicit runbooks.",
            "mitigation": "Schedule failover validation before production cutover.",
        })
    do_nothing = [
        {
            "risk": "Technical debt and platform constraints continue to accumulate.",
            "impact": "Higher delivery risk and slower response to business demand.",
        },
        {
            "risk": "Commercial and capacity assumptions remain unvalidated.",
            "impact": "Procurement or renewal decisions may be made without a current baseline.",
        },
    ]
    return {
        "migration_risks": migration[:4],
        "do_nothing_risks": do_nothing[:4],
    }


def _infer_business_case_roadmap(text: str) -> dict:
    total_weeks = _timeline_to_weeks(text)
    if total_weeks and total_weeks >= 6:
        define = max(1, round(total_weeks * 0.25))
        design = max(1, round(total_weeks * 0.35))
        deliver = max(1, total_weeks - define - design)
        total_duration = f"{total_weeks} weeks"
    else:
        define, design, deliver = 2, 4, 6
        total_duration = "12 weeks (to validate)"
    return {
        "total_duration": total_duration,
        "phases": [
            {
                "name": "Define",
                "duration": f"{define} weeks",
                "deliverables": ["Scope baseline", "Business case validation"],
                "milestones": ["Stakeholder alignment"],
            },
            {
                "name": "Design",
                "duration": f"{design} weeks",
                "deliverables": ["Target architecture", "Security and operating model"],
                "milestones": ["Design sign-off"],
            },
            {
                "name": "Deliver",
                "duration": f"{deliver} weeks",
                "deliverables": ["Migration execution", "Cutover and stabilization"],
                "milestones": ["Go-live readiness"],
            },
        ],
    }


def _infer_business_case_recommendation(text: str) -> dict:
    objective = _infer_driver_objective(text)
    return {
        "summary": f"Proceed to detailed design on OCI to validate {objective} with customer-specific sizing and commercials.",
        "next_steps": [
            "Confirm scope, current-state baseline, and success criteria.",
            "Validate commercials, licensing assumptions, and target delivery model.",
            "Run a design workshop and prepare the first migration wave.",
        ],
    }


def _enrich_sparse_business_case(spec: dict) -> dict:
    """Fill safe narrative sections so sparse business cases still render professionally."""
    if not isinstance(spec, dict):
        return spec
    bc = spec.get("business_case", spec)
    if not isinstance(bc, dict):
        return spec

    text = _extract_business_case_text(bc)
    if not text.strip():
        return spec

    enriched_bc = dict(bc)
    if "drivers" not in enriched_bc:
        enriched_bc["drivers"] = _infer_business_case_drivers(text)
    if "risks" not in enriched_bc:
        enriched_bc["risks"] = _infer_business_case_risks(text)
    if "roadmap" not in enriched_bc:
        enriched_bc["roadmap"] = _infer_business_case_roadmap(text)
    if "recommendation" not in enriched_bc:
        enriched_bc["recommendation"] = _infer_business_case_recommendation(text)

    if "business_case" in spec:
        enriched = dict(spec)
        enriched["business_case"] = enriched_bc
        return enriched
    return enriched_bc


# ============================================================
# Business Case Deck Generator
# ============================================================

class BusinessCaseDeckGenerator(OraclePresBase):
    """Generate business case decks using Oracle FY26 template layouts."""

    # Use Teal accent for business case slides (finance / CxO audience)
    TITLE_ACCENT_COLOR = Colors.TEAL

    # ================================================================
    # Slide Methods
    # ================================================================

    def add_cover_slide(self, customer: str, subtitle: str = "",
                        prepared_by: str = "", date: str = "",
                        prepared_by_role: str = ""):
        """Slide 1: Cover using Dark Title_Pillar layout."""
        slide = self._add_layout_slide(Layouts.COVER_DARK)
        self._set_placeholder(slide, 0, customer)  # Title
        self._set_placeholder(slide, 33, subtitle or "Business Case for Oracle Cloud Infrastructure")
        date_str = date or datetime.now().strftime("%B %Y")
        self._set_placeholder(slide, 35, date_str)
        if prepared_by:
            byline = f"{prepared_by}, {prepared_by_role}" if prepared_by_role else prepared_by
            self._set_placeholder(slide, 34, f"Prepared by: {byline}")

    def add_executive_summary_slide(self, statement: str):
        """Slide 2: Executive Summary — controlled typography on blank slide."""
        slide = self._add_blank_slide()

        self._add_title_bar(slide, "Executive Summary")

        # Body — large enough to read, small enough to fit the paragraph
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = statement
        p.font.size = Pt(18)
        p.font.name = self.FONT
        p.font.color.rgb = Colors.PRIMARY_TEXT
        p.space_after = Pt(6)

    def add_business_drivers_slide(self, drivers: list):
        """Slide 3: Business Drivers — numbered cards on blank slide.

        drivers: list of up to 3 strings, each a driver statement (may include newline
        separating a bold headline from detail text).
        """
        slide = self._add_blank_slide()

        self._add_title_bar(slide, "Business Drivers")

        card_colors = [Colors.TEAL, Colors.BURNT_ORANGE, Colors.FOREST]
        card_y_start = Inches(1.3)
        available_h = Inches(5.8)
        n = min(len(drivers), 3)
        card_h = available_h / n if n else available_h
        card_h = min(card_h, Inches(2.1))

        for i, driver in enumerate(drivers[:3]):
            color = card_colors[i % len(card_colors)]
            y = card_y_start + i * card_h

            # Left accent bar
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y + Inches(0.1), Inches(0.08), card_h - Inches(0.2))
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            # Number badge
            self._add_textbox(
                slide, Inches(1.05), y + Inches(0.12), Inches(0.5), Inches(0.45),
                text=f"0{i+1}", font_size=18, bold=True, color=color,
            )

            # Split driver into headline + detail (split on \n if present)
            parts = driver.split("\n", 1)
            headline = parts[0].strip()
            detail = parts[1].strip() if len(parts) > 1 else ""

            # Headline
            self._add_textbox(
                slide, Inches(1.6), y + Inches(0.1), Inches(11.0), Inches(0.5),
                text=headline, font_size=15, bold=True, color=Colors.PRIMARY_TEXT,
            )
            # Detail
            if detail:
                self._add_textbox(
                    slide, Inches(1.6), y + Inches(0.6), Inches(11.0), card_h - Inches(0.75),
                    text=detail, font_size=13, color=Colors.SECONDARY_TEXT,
                )

    def add_tco_slide(self, tco: dict):
        """Slide 4: TCO Comparison — table on blank slide.

        tco: dict from business-case.yaml with current_state and proposed_oci.
        """
        slide = self._add_blank_slide()

        self._add_title_bar(slide, "Total Cost of Ownership")

        # Subtitle
        horizon = tco.get("horizon_years", 3)
        self._add_textbox(
            slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
            text=tco.get("comparison_label", f"{horizon}-Year Comparison  |  Current State vs Oracle Cloud Infrastructure"),
            font_size=13, color=Colors.TEAL, bold=True,
        )

        current = tco.get("current_state", {})
        proposed = tco.get("proposed_oci", {})
        savings = tco.get("savings", {})

        ann_current = current.get("total_annual", 0)
        ann_oci = proposed.get("total_annual", 0)
        ann_savings = savings.get("annual", 0) or (ann_current - ann_oci)
        monthly_current = ann_current / 12 if ann_current else 0
        monthly_oci = ann_oci / 12 if ann_oci else 0

        kpi_labels = tco.get("kpi_labels", {}) if isinstance(tco.get("kpi_labels", {}), dict) else {}
        kpis = [
            (f"${monthly_current:,.0f}/mo", kpi_labels.get("current", "Current state")),
            (f"${monthly_oci:,.0f}/mo", kpi_labels.get("proposed", "Proposed OCI")),
            (f"${abs(ann_savings) / 12:,.0f}/mo", kpi_labels.get("delta", "Monthly delta")),
        ]
        for idx, (value, label) in enumerate(kpis):
            x = Inches(0.8 + idx * 4.0)
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.55), Inches(3.65), Inches(1.0))
            box.fill.solid()
            box.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
            box.line.color.rgb = Colors.MUTED_TEAL
            self._add_textbox(slide, x + Inches(0.15), Inches(1.68), Inches(3.35), Inches(0.4),
                              text=value, font_size=20, bold=True, color=Colors.TEAL,
                              alignment=PP_ALIGN.CENTER)
            self._add_textbox(slide, x + Inches(0.15), Inches(2.08), Inches(3.35), Inches(0.28),
                              text=label, font_size=10, color=Colors.SECONDARY_TEXT,
                              alignment=PP_ALIGN.CENTER)

        row_labels = tco.get("row_labels", {}) if isinstance(tco.get("row_labels", {}), dict) else {}
        rows_data = [
            (row_labels.get("infrastructure", "Infrastructure"), current.get("annual_infrastructure", 0), proposed.get("annual_cloud_consumption", 0)),
            ("Licensing / Support", current.get("annual_licensing", 0), proposed.get("annual_licensing", 0)),
            ("Operations (People)", current.get("annual_operations", 0), proposed.get("annual_operations", 0)),
            ("Downtime Cost", current.get("annual_downtime_cost", 0), proposed.get("annual_downtime_cost", 0)),
            ("Compliance", current.get("annual_compliance_cost", 0), proposed.get("annual_compliance_cost", 0)),
        ]
        # Filter out zero rows
        rows_data = [(n, c, o) for n, c, o in rows_data if c or o]

        # Add migration one-time
        migration = proposed.get("migration_one_time", 0)

        num_rows = len(rows_data) + 3  # header + rows + annual total + horizon total
        table_top = Inches(3.0)
        table = self._add_table(
            slide, num_rows, 4,
            Inches(0.8), table_top,
            Inches(11.7), Inches(0.42 * num_rows),
        )
        table.columns[0].width = Inches(3.5)
        table.columns[1].width = Inches(2.7)
        table.columns[2].width = Inches(2.7)
        table.columns[3].width = Inches(2.8)

        # Header
        column_labels = tco.get("column_labels", {}) if isinstance(tco.get("column_labels", {}), dict) else {}
        headers = [
            "Cost Category",
            column_labels.get("current", "Current (Annual)"),
            column_labels.get("proposed", "OCI (Annual)"),
            column_labels.get("delta", "Savings"),
        ]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
            )

        # Data rows
        for i, (name, curr, oci) in enumerate(rows_data):
            row = i + 1
            bg = Colors.TABLE_ALT_ROW if row % 2 == 0 else None
            curr_val = curr if isinstance(curr, (int, float)) else 0
            oci_val = oci if isinstance(oci, (int, float)) else 0
            save_val = curr_val - oci_val
            self._style_table_cell(table.cell(row, 0), name, font_size=10, bg_color=bg)
            self._style_table_cell(table.cell(row, 1), f"${curr_val:,.0f}", font_size=10, bg_color=bg, alignment=PP_ALIGN.RIGHT)
            self._style_table_cell(table.cell(row, 2), f"${oci_val:,.0f}", font_size=10, bg_color=bg, alignment=PP_ALIGN.RIGHT)
            save_color = Colors.SUCCESS if save_val > 0 else Colors.ERROR
            self._style_table_cell(table.cell(row, 3), f"${save_val:,.0f}", font_size=10, bold=True, color=save_color, bg_color=bg, alignment=PP_ALIGN.RIGHT)

        # Annual total row
        total_row = len(rows_data) + 1
        ann_current = ann_current or sum(r[1] for r in rows_data if isinstance(r[1], (int, float)))
        ann_oci = ann_oci or sum(r[2] for r in rows_data if isinstance(r[2], (int, float)))
        ann_savings = ann_savings or (ann_current - ann_oci)
        self._style_table_cell(table.cell(total_row, 0), "TOTAL ANNUAL", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL)
        self._style_table_cell(table.cell(total_row, 1), f"${ann_current:,.0f}", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL, alignment=PP_ALIGN.RIGHT)
        self._style_table_cell(table.cell(total_row, 2), f"${ann_oci:,.0f}", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL, alignment=PP_ALIGN.RIGHT)
        self._style_table_cell(table.cell(total_row, 3), f"${ann_savings:,.0f}", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL, alignment=PP_ALIGN.RIGHT)

        # Horizon total row
        h_row = total_row + 1
        h_current = current.get("total_over_horizon", 0) or (ann_current * horizon)
        h_oci = proposed.get("total_over_horizon", 0) or (ann_oci * horizon + migration)
        h_savings = savings.get("over_horizon", 0) or (h_current - h_oci)
        pct = savings.get("percentage", 0) or (h_savings / h_current * 100 if h_current else 0)
        self._style_table_cell(table.cell(h_row, 0), f"TOTAL {horizon}-YEAR", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.DARK_BG)
        self._style_table_cell(table.cell(h_row, 1), f"${h_current:,.0f}", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.DARK_BG, alignment=PP_ALIGN.RIGHT)
        self._style_table_cell(table.cell(h_row, 2), f"${h_oci:,.0f}", font_size=11, bold=True, color=Colors.WHITE, bg_color=Colors.DARK_BG, alignment=PP_ALIGN.RIGHT)
        self._style_table_cell(table.cell(h_row, 3), f"${h_savings:,.0f} ({pct:.0f}%)", font_size=11, bold=True, color=Colors.GOLD, bg_color=Colors.DARK_BG, alignment=PP_ALIGN.RIGHT)

        # Migration note
        if migration:
            note_y = table_top + Inches(0.42 * num_rows) + Inches(0.15)
            self._add_textbox(
                slide, Inches(0.8), note_y, Inches(11), Inches(0.3),
                text=f"* Includes one-time migration investment of ${migration:,.0f} in Year 1",
                font_size=9, italic=True, color=Colors.SECONDARY_TEXT,
            )

        # Assumptions
        assumptions = tco.get("assumptions", [])
        if assumptions:
            a_y = table_top + Inches(0.42 * num_rows) + Inches(0.45)
            self._add_textbox(
                slide, Inches(0.8), a_y, Inches(11), Inches(0.25),
                text="Assumptions:", font_size=9, bold=True, color=Colors.SECONDARY_TEXT,
            )
            for idx, a in enumerate(assumptions[:7]):
                self._add_textbox(
                    slide, Inches(1.0), a_y + Inches(0.25 + idx * 0.22),
                    Inches(11), Inches(0.22),
                    text=f"• {a}", font_size=7, italic=True, color=Colors.SECONDARY_TEXT,
                )

    def add_roi_slide(self, roi: dict, headline: str = ""):
        """Slide 5: ROI — centered big number + 3 supporting metrics."""
        slide = self._add_blank_slide()

        self._add_title_bar(slide, "Return on Investment")

        cards = roi.get("cards", [])
        if cards:
            metric_text = headline or pick(roi, "headline", default="ROI logic")
            self._add_textbox(
                slide, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.62),
                text=metric_text, font_size=34, bold=True, color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER,
            )
            summary = pick(roi, "label", "summary")
            if summary:
                self._add_textbox(
                    slide, Inches(1.0), Inches(1.82), Inches(11.3), Inches(0.5),
                    text=summary, font_size=13, color=Colors.SECONDARY_TEXT,
                    alignment=PP_ALIGN.CENTER,
                )

            card_w = Inches(3.65)
            card_h = Inches(2.25)
            gap = Inches(0.35)
            x0 = Inches(0.9)
            y0 = Inches(2.75)
            colors = [Colors.TEAL, Colors.BURNT_ORANGE, Colors.FOREST]
            for idx, card in enumerate(cards[:3]):
                x = x0 + idx * (card_w + gap)
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, card_w, card_h)
                bg.fill.solid()
                bg.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
                bg.line.color.rgb = colors[idx % len(colors)]
                self._add_textbox(
                    slide, x + Inches(0.18), y0 + Inches(0.18), card_w - Inches(0.36), Inches(0.28),
                    text=card.get("title", ""), font_size=11, bold=True, color=colors[idx % len(colors)],
                )
                self._add_textbox(
                    slide, x + Inches(0.18), y0 + Inches(0.58), card_w - Inches(0.36), Inches(0.48),
                    text=card.get("metric", ""), font_size=18, bold=True, color=Colors.PRIMARY_TEXT,
                )
                self._add_textbox(
                    slide, x + Inches(0.18), y0 + Inches(1.14), card_w - Inches(0.36), Inches(0.88),
                    text=card.get("detail", ""), font_size=8, color=Colors.SECONDARY_TEXT,
                )

            note = pick(roi, "note")
            if note:
                self._add_textbox(
                    slide, Inches(0.9), Inches(5.45), Inches(11.5), Inches(0.65),
                    text=note, font_size=10, italic=True, color=Colors.SECONDARY_TEXT,
                    alignment=PP_ALIGN.CENTER,
                )
            return

        # Big centered ROI number
        pct = roi.get("three_year_roi_pct", 0)
        metric_text = headline or (f"{pct:.0f}%" if pct else "—")
        big_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.4))
        big_box.text_frame.word_wrap = False
        p = big_box.text_frame.paragraphs[0]
        p.text = metric_text
        p.font.size = Pt(96)
        p.font.bold = True
        p.font.name = self.FONT_HEADING
        p.font.color.rgb = Colors.TEAL
        p.alignment = PP_ALIGN.CENTER

        # Label under big number
        self._add_textbox(
            slide, Inches(0.8), Inches(3.75), Inches(11.7), Inches(0.4),
            text=roi.get("label", "3-Year Return on Investment"), font_size=16,
            color=Colors.SECONDARY_TEXT, alignment=PP_ALIGN.CENTER,
        )

        # 3 supporting metric boxes
        payback = roi.get("payback_months", 0)
        investment = roi.get("total_investment", 0)
        benefit = roi.get("annual_net_benefit", 0)
        metrics = []
        if payback:
            metrics.append((f"{payback} months", "Payback Period"))
        if investment:
            metrics.append((f"${investment:,.0f}", "Total Investment"))
        if benefit:
            metrics.append((f"${benefit:,.0f}/yr", "Annual Net Benefit"))

        box_w = Inches(3.5)
        gap = Inches(0.45)
        total_w = box_w * len(metrics) + gap * (len(metrics) - 1)
        x_start = (Inches(13.33) - total_w) / 2
        y_box = Inches(4.4)

        for j, (value, label) in enumerate(metrics):
            x = x_start + j * (box_w + gap)
            # Background
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_box, box_w, Inches(1.4))
            bg.fill.solid()
            bg.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
            bg.line.color.rgb = Colors.MUTED_TEAL

            # Value
            self._add_textbox(
                slide, x, y_box + Inches(0.15), box_w, Inches(0.7),
                text=value, font_size=26, bold=True, color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER,
            )
            # Label
            self._add_textbox(
                slide, x, y_box + Inches(0.85), box_w, Inches(0.45),
                text=label, font_size=11, color=Colors.SECONDARY_TEXT,
                alignment=PP_ALIGN.CENTER,
            )

    def add_tco_projection_slide(
        self,
        projection: list,
        title: str = "3-Year TCO Projection",
        storage_economics: dict | None = None,
    ):
        """Slide: multi-year TCO forecast with CPU and storage assumptions."""
        if not projection:
            return
        slide = self._add_blank_slide()
        self._add_title_bar(slide, title)

        headers = ["Period", "CPU Demand", "Storage", "ADB-S As-Is", "ADB-D To-Be", "Delta", "Note"]
        rows = projection[:4]
        table = self._add_table(
            slide, len(rows) + 1, len(headers),
            Inches(0.55), Inches(1.35),
            Inches(12.25), Inches(0.55 * (len(rows) + 1)),
        )
        widths = [1.05, 1.65, 1.8, 1.65, 1.65, 1.35, 3.1]
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)

        for j, header in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), header, font_size=9, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT,
            )

        for i, row in enumerate(rows, start=1):
            bg = Colors.TABLE_ALT_ROW if i % 2 == 0 else None
            values = [
                row.get("period", ""),
                row.get("cpu", ""),
                row.get("storage", ""),
                row.get("as_is", ""),
                row.get("to_be", ""),
                row.get("delta", ""),
                row.get("note", ""),
            ]
            for j, value in enumerate(values):
                self._style_table_cell(
                    table.cell(i, j), value, font_size=8, bg_color=bg,
                    alignment=PP_ALIGN.RIGHT if j in (3, 4, 5) else PP_ALIGN.LEFT,
                    color=Colors.ERROR if j == 5 and str(value).startswith("+") else Colors.PRIMARY_TEXT,
                )

        footnote = rows[0].get("footnote", "") if isinstance(rows[0], dict) else ""
        if footnote:
            self._add_textbox(
                slide, Inches(0.65), Inches(6.25), Inches(12.0), Inches(0.45),
                text=footnote, font_size=8, italic=True, color=Colors.SECONDARY_TEXT,
            )

        if isinstance(storage_economics, dict) and storage_economics:
            y = Inches(4.35)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), y, Inches(11.85), Inches(1.55))
            box.fill.solid()
            box.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
            box.line.color.rgb = Colors.MUTED_TEAL

            self._add_textbox(
                slide, Inches(0.95), y + Inches(0.12), Inches(11.45), Inches(0.26),
                text=pick(storage_economics, "headline", default="Storage economics"),
                font_size=10, bold=True, color=Colors.TEAL,
            )

            cards = storage_economics.get("cards", [])
            if cards:
                card_w = Inches(3.55)
                for idx, card in enumerate(cards[:3]):
                    x = Inches(0.95) + idx * Inches(3.8)
                    self._add_textbox(
                        slide, x, y + Inches(0.45), card_w, Inches(0.35),
                        text=card.get("value", ""), font_size=15, bold=True, color=Colors.PRIMARY_TEXT,
                    )
                    self._add_textbox(
                        slide, x, y + Inches(0.82), card_w, Inches(0.34),
                        text=card.get("label", ""), font_size=8, color=Colors.SECONDARY_TEXT,
                    )
                    detail = card.get("detail", "")
                    if detail:
                        self._add_textbox(
                            slide, x, y + Inches(1.13), card_w, Inches(0.28),
                            text=detail, font_size=7, color=Colors.SECONDARY_TEXT,
                        )

    def add_tco_crossover_slide(self, chart_spec: dict):
        """Slide: show when ADB-D becomes cheaper than ADB-S."""
        if not isinstance(chart_spec, dict):
            return

        categories = chart_spec.get("categories", [])
        as_is = chart_spec.get("as_is", [])
        to_be = chart_spec.get("to_be", [])
        if not categories or not as_is or not to_be:
            return

        slide = self._add_blank_slide()
        self._add_title_bar(slide, pick(chart_spec, "title", default="TCO Crossover"))

        subtitle = pick(chart_spec, "subtitle")
        if subtitle:
            self._add_textbox(
                slide, Inches(0.75), Inches(1.05), Inches(12.0), Inches(0.35),
                text=subtitle, font_size=13, color=Colors.SECONDARY_TEXT,
            )

        # Draw the chart with native shapes so it remains visible in lightweight
        # renderers that do not paint embedded Office chart objects.
        y_min = float(chart_spec.get("y_axis_min", min(as_is + to_be) * 0.9))
        y_max = float(chart_spec.get("y_axis_max", max(as_is + to_be) * 1.1))
        y_step = float(chart_spec.get("y_axis_major_unit", 0.5))
        plot_x, plot_y = Inches(0.95), Inches(1.8)
        plot_w, plot_h = Inches(7.75), Inches(4.2)
        axis_color = Colors.SECONDARY_TEXT
        grid_color = RGBColor(220, 218, 214)
        series_colors = [Colors.BURNT_ORANGE, Colors.TEAL]

        def x_pos(idx: int) -> int:
            if len(categories) == 1:
                return int(plot_x + plot_w / 2)
            return int(plot_x + idx * (plot_w / (len(categories) - 1)))

        def y_pos(value: float) -> int:
            ratio = (float(value) - y_min) / (y_max - y_min)
            return int(plot_y + plot_h - ratio * plot_h)

        # Grid and y labels.
        tick = y_min
        while tick <= y_max + 0.001:
            y = y_pos(tick)
            grid = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, plot_x, y, plot_x + plot_w, y)
            grid.line.color.rgb = grid_color
            grid.line.width = Pt(0.6)
            self._add_textbox(
                slide, Inches(0.52), y - Inches(0.1), Inches(0.35), Inches(0.2),
                text=f"${tick:.1f}M", font_size=7, color=Colors.SECONDARY_TEXT,
                alignment=PP_ALIGN.RIGHT,
            )
            tick += y_step

        # Axes.
        x_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, plot_x, plot_y + plot_h, plot_x + plot_w, plot_y + plot_h)
        x_axis.line.color.rgb = axis_color
        y_axis = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, plot_x, plot_y, plot_x, plot_y + plot_h)
        y_axis.line.color.rgb = axis_color

        baseline_y = y_pos(y_min)
        group_w = plot_w / len(categories)
        bar_w = min(Inches(0.28), group_w * 0.28)
        gap_w = Inches(0.06)

        for idx, category in enumerate(categories):
            center = int(plot_x + group_w * (idx + 0.5))
            values = [as_is[idx], to_be[idx]]
            for s_idx, value in enumerate(values):
                color = series_colors[s_idx]
                x = center - bar_w - gap_w / 2 if s_idx == 0 else center + gap_w / 2
                y = y_pos(value)
                h = baseline_y - y
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(bar_w), int(h))
                bar.fill.solid()
                bar.fill.fore_color.rgb = color
                bar.line.fill.background()
                self._add_textbox(
                    slide, int(x) - Inches(0.16), y - Inches(0.26), bar_w + Inches(0.32), Inches(0.2),
                    text=f"${value:.2f}M", font_size=7, color=color,
                    alignment=PP_ALIGN.CENTER,
                )

        for idx, category in enumerate(categories):
            center = int(plot_x + group_w * (idx + 0.5))
            self._add_textbox(
                slide, center - Inches(0.45), plot_y + plot_h + Inches(0.1),
                Inches(0.9), Inches(0.25), text=str(category), font_size=8,
                color=Colors.SECONDARY_TEXT, alignment=PP_ALIGN.CENTER,
            )

        # Legend.
        legend_y = plot_y + plot_h + Inches(0.55)
        for idx, (label, color) in enumerate([
            (pick(chart_spec, "as_is_label", default="ADB-S"), series_colors[0]),
            (pick(chart_spec, "to_be_label", default="ADB-D"), series_colors[1]),
        ]):
            x = plot_x + Inches(idx * 2.2)
            swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, legend_y, Inches(0.18), Inches(0.08))
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = color
            swatch.line.fill.background()
            self._add_textbox(slide, x + Inches(0.25), legend_y - Inches(0.06), Inches(1.6), Inches(0.22), text=label, font_size=8, color=Colors.SECONDARY_TEXT)

        callout = pick(chart_spec, "callout", default="")
        if callout:
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(9.15), Inches(1.7), Inches(3.35), Inches(1.25),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = Colors.TEAL
            box.line.fill.background()
            self._add_textbox(
                slide, Inches(9.35), Inches(1.92), Inches(2.95), Inches(0.78),
                text=callout, font_size=17, bold=True, color=Colors.WHITE,
                alignment=PP_ALIGN.CENTER,
            )

        bullets = chart_spec.get("bullets", [])
        y = Inches(3.25)
        for idx, bullet in enumerate(bullets[:4]):
            self._add_textbox(
                slide, Inches(9.2), y + idx * Inches(0.55), Inches(3.25), Inches(0.42),
                text=f"- {bullet}", font_size=9, color=Colors.SECONDARY_TEXT,
            )

        note = pick(chart_spec, "note")
        if note:
            self._add_textbox(
                slide, Inches(0.85), Inches(6.35), Inches(11.8), Inches(0.35),
                text=note, font_size=8, italic=True, color=Colors.SECONDARY_TEXT,
            )

    def add_cost_breakdown_slide(self, breakdown: dict):
        """Slide: explain how cloud services and operations costs are built."""
        if not isinstance(breakdown, dict):
            return
        slide = self._add_blank_slide()
        self._add_title_bar(slide, pick(breakdown, "title", default="BOM + Operations Cost Breakdown"))

        scenarios = breakdown.get("scenarios", [])
        if not scenarios:
            return

        col_w = Inches(5.8)
        gap = Inches(0.55)
        x_positions = [Inches(0.75), Inches(0.75) + col_w + gap]
        colors = [Colors.TEAL, Colors.BURNT_ORANGE]

        for idx, scenario in enumerate(scenarios[:2]):
            x = x_positions[idx]
            y = Inches(1.3)
            color = colors[idx % len(colors)]
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_w, Inches(0.45))
            header.fill.solid()
            header.fill.fore_color.rgb = color
            header.line.fill.background()
            self._add_textbox(
                slide, x + Inches(0.15), y + Inches(0.06), col_w - Inches(0.3), Inches(0.32),
                text=pick(scenario, "name", default=f"Scenario {idx+1}"),
                font_size=12, bold=True, color=Colors.WHITE,
            )

            rows = scenario.get("lines", [])
            table = self._add_table(
                slide, len(rows) + 2, 3,
                x, y + Inches(0.6),
                col_w, Inches(0.32 * (len(rows) + 2)),
            )
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(1.35)
            table.columns[2].width = Inches(1.45)
            for j, h in enumerate(["Cost Item", "Monthly", "Annual"]):
                self._style_table_cell(
                    table.cell(0, j), h, font_size=8, bold=True,
                    color=Colors.WHITE, bg_color=color,
                    alignment=PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT,
                )
            for r_idx, row in enumerate(rows, start=1):
                bg = Colors.TABLE_ALT_ROW if r_idx % 2 == 0 else None
                self._style_table_cell(table.cell(r_idx, 0), row.get("item", ""), font_size=7, bg_color=bg)
                self._style_table_cell(table.cell(r_idx, 1), row.get("monthly", ""), font_size=7, bg_color=bg, alignment=PP_ALIGN.RIGHT)
                self._style_table_cell(table.cell(r_idx, 2), row.get("annual", ""), font_size=7, bg_color=bg, alignment=PP_ALIGN.RIGHT)
            total = scenario.get("total", {})
            last = len(rows) + 1
            self._style_table_cell(table.cell(last, 0), "TOTAL", font_size=8, bold=True, color=Colors.WHITE, bg_color=Colors.DARK_BG)
            self._style_table_cell(table.cell(last, 1), total.get("monthly", ""), font_size=8, bold=True, color=Colors.WHITE, bg_color=Colors.DARK_BG, alignment=PP_ALIGN.RIGHT)
            self._style_table_cell(table.cell(last, 2), total.get("annual", ""), font_size=8, bold=True, color=Colors.WHITE, bg_color=Colors.DARK_BG, alignment=PP_ALIGN.RIGHT)

        notes = breakdown.get("notes", [])
        if notes:
            y = Inches(6.0)
            self._add_textbox(slide, Inches(0.8), y, Inches(11.8), Inches(0.25), text="Assumptions:", font_size=8, bold=True, color=Colors.SECONDARY_TEXT)
            for idx, note in enumerate(notes[:4]):
                self._add_textbox(
                    slide, Inches(1.0), y + Inches(0.23 + idx * 0.2), Inches(11.6), Inches(0.2),
                    text=f"- {note}", font_size=7, color=Colors.SECONDARY_TEXT,
                )

    def add_business_value_slide(self, value_case: dict):
        """Slide: risk-adjusted value / avoided downtime model."""
        if not isinstance(value_case, dict):
            return
        slide = self._add_blank_slide()
        self._add_title_bar(slide, pick(value_case, "title", default="Risk-Adjusted Business Value"))

        headline = pick(value_case, "headline")
        if headline:
            self._add_textbox(
                slide, Inches(0.8), Inches(1.15), Inches(11.8), Inches(0.5),
                text=headline, font_size=16, bold=True, color=Colors.TEAL,
            )

        rows = value_case.get("rows", [])
        if rows:
            headers = ["Value Area", "How to Measure", "Business Case Treatment"]
            table = self._add_table(slide, len(rows) + 1, 3, Inches(0.75), Inches(1.85), Inches(12.0), Inches(0.75 * (len(rows) + 1)))
            table.columns[0].width = Inches(2.6)
            table.columns[1].width = Inches(4.3)
            table.columns[2].width = Inches(5.1)
            for j, h in enumerate(headers):
                self._style_table_cell(table.cell(0, j), h, font_size=9, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL)
            for i, row in enumerate(rows, start=1):
                bg = Colors.TABLE_ALT_ROW if i % 2 == 0 else None
                self._style_table_cell(table.cell(i, 0), row.get("area", ""), font_size=8, bg_color=bg)
                self._style_table_cell(table.cell(i, 1), row.get("measure", ""), font_size=8, bg_color=bg)
                self._style_table_cell(table.cell(i, 2), row.get("treatment", ""), font_size=8, bg_color=bg)

        break_even = value_case.get("break_even", [])
        if break_even:
            self._add_textbox(
                slide, Inches(0.8), Inches(5.25), Inches(11.8), Inches(0.3),
                text="Break-even avoided outage impact", font_size=10, bold=True, color=Colors.SECONDARY_TEXT,
            )
            headers = ["Horizon", "Incremental TCO", "If avoiding 4h/year", "If avoiding 8h/year"]
            table = self._add_table(slide, len(break_even) + 1, 4, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.38 * (len(break_even) + 1)))
            for j, h in enumerate(headers):
                self._style_table_cell(table.cell(0, j), h, font_size=8, bold=True, color=Colors.WHITE, bg_color=Colors.BURNT_ORANGE)
            for i, row in enumerate(break_even, start=1):
                bg = Colors.TABLE_ALT_ROW if i % 2 == 0 else None
                values = [row.get("horizon", ""), row.get("investment", ""), row.get("four_hours", ""), row.get("eight_hours", "")]
                for j, value in enumerate(values):
                    self._style_table_cell(table.cell(i, j), value, font_size=8, bg_color=bg, alignment=PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT)

    def add_value_drivers_slide(self, drivers: list):
        """Slide 6: Value Drivers — 4 categories on blank slide.

        drivers: list of {"category": str, "title": str, "description": str, "quantified": str}
        """
        drivers = [
            {"title": driver, "description": ""} if isinstance(driver, str) else driver
            for driver in drivers
        ]
        slide = self._add_blank_slide()

        self._add_title_bar(slide, "Value Drivers")

        # 4 value cards in a 2x2 grid
        card_colors = [Colors.TEAL, Colors.FOREST, Colors.BURNT_ORANGE, Colors.ORACLE_RED]
        positions = [
            (Inches(0.8), Inches(1.5)),   # top-left
            (Inches(6.8), Inches(1.5)),   # top-right
            (Inches(0.8), Inches(4.3)),   # bottom-left
            (Inches(6.8), Inches(4.3)),   # bottom-right
        ]
        card_w, card_h = Inches(5.5), Inches(2.5)

        for i, driver in enumerate(drivers[:4]):
            x, y = positions[i]
            color = card_colors[i % len(card_colors)]

            # Color accent bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()

            # Title
            title = driver.get("title", driver.get("category", "").replace("_", " ").title())
            self._add_textbox(
                slide, x + Inches(0.25), y + Inches(0.1),
                card_w - Inches(0.3), Inches(0.4),
                text=title, font_size=16, bold=True, color=color,
            )

            # Quantified value (big)
            quantified = driver.get("quantified", "")
            if quantified:
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(0.55),
                    card_w - Inches(0.3), Inches(0.5),
                    text=quantified, font_size=22, bold=True, color=Colors.PRIMARY_TEXT,
                )

            # Description
            desc = driver.get("description", "")
            if desc:
                self._add_textbox(
                    slide, x + Inches(0.25), y + Inches(1.2),
                    card_w - Inches(0.3), Inches(1.0),
                    text=desc, font_size=11, color=Colors.SECONDARY_TEXT,
                )

    def add_risk_slide(self, migration_risks: list, do_nothing_risks: list):
        """Slide 7: Risk Assessment — two-column table layout."""
        slide = self._add_blank_slide()

        self._add_title_bar(slide, "Risk Assessment")

        col_w = Inches(5.6)
        col_gap = Inches(0.5)
        left_x = Inches(0.8)
        right_x = left_x + col_w + col_gap
        header_y = Inches(1.2)
        content_start_y = Inches(1.75)
        row_h = Inches(0.8)          # height per risk block (title + detail)
        slide_h = Inches(7.1)

        # Stretch rows to fill available slide height
        n_rows = max(len(migration_risks[:5]), len(do_nothing_risks[:5]))
        if n_rows > 0:
            available = slide_h - content_start_y
            row_h = min(Inches(1.8), available / n_rows)

        # Column header backgrounds
        for x, color, label in [
            (left_x,  Colors.TEAL,  "Migration Risks (Mitigated)"),
            (right_x, Colors.ERROR, "Risks of Inaction"),
        ]:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, header_y, col_w, Inches(0.45))
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()
            self._add_textbox(
                slide, x + Inches(0.15), header_y + Inches(0.05), col_w - Inches(0.2), Inches(0.38),
                text=label, font_size=13, bold=True, color=Colors.WHITE,
            )

        # Left column: migration risks
        for i, risk in enumerate(migration_risks[:5]):
            r_text = risk.get("risk", str(risk)) if isinstance(risk, dict) else str(risk)
            mitigation = risk.get("mitigation", "") if isinstance(risk, dict) else ""
            y = content_start_y + i * row_h

            # Alternating row background
            if i % 2 == 0:
                row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, y, col_w, row_h)
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
                row_bg.line.fill.background()

            self._add_textbox(
                slide, left_x + Inches(0.15), y + Inches(0.07), col_w - Inches(0.25), Inches(0.4),
                text=f"- {r_text}", font_size=12, bold=True, color=Colors.PRIMARY_TEXT,
            )
            if mitigation:
                self._add_textbox(
                    slide, left_x + Inches(0.25), y + Inches(0.45), col_w - Inches(0.35), row_h - Inches(0.48),
                    text=f"Mitigation: {mitigation}", font_size=11,
                    italic=True, color=Colors.SECONDARY_TEXT,
                )

        # Right column: do-nothing risks
        for i, risk in enumerate(do_nothing_risks[:5]):
            r_text = risk.get("risk", str(risk)) if isinstance(risk, dict) else str(risk)
            impact = risk.get("impact", "") if isinstance(risk, dict) else ""
            timeline = risk.get("timeline", "") if isinstance(risk, dict) else ""
            y = content_start_y + i * row_h

            if i % 2 == 0:
                row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, y, col_w, row_h)
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xF3)
                row_bg.line.fill.background()

            self._add_textbox(
                slide, right_x + Inches(0.15), y + Inches(0.07), col_w - Inches(0.25), Inches(0.4),
                text=f"- {r_text}", font_size=12, bold=True, color=Colors.ERROR,
            )
            detail_parts = []
            if impact:
                detail_parts.append(f"Impact: {impact}")
            if timeline:
                detail_parts.append(f"By: {timeline}")
            if detail_parts:
                self._add_textbox(
                    slide, right_x + Inches(0.25), y + Inches(0.45), col_w - Inches(0.35), row_h - Inches(0.48),
                    text="  |  ".join(detail_parts), font_size=10,
                    italic=True, color=Colors.SECONDARY_TEXT,
                )

    def add_roadmap_slide(self, phases: list, total_duration: str = ""):
        """Slide 8: Implementation Roadmap — timeline bars filling the slide."""
        slide = self._add_blank_slide()

        title = "Implementation Roadmap"
        if total_duration:
            title += f"  —  {total_duration}"
        self._add_title_bar(slide, title)

        phase_colors = [Colors.TEAL, Colors.BURNT_ORANGE, Colors.FOREST, Colors.ORACLE_RED]
        content_start = Inches(1.25)
        available_h = Inches(5.9)
        n = min(len(phases), 4)
        row_h = available_h / n if n else available_h

        bar_left = Inches(3.6)
        bar_w = Inches(9.1)
        label_w = Inches(2.6)

        for i, phase in enumerate(phases[:4]):
            color = phase_colors[i % len(phase_colors)]
            y = content_start + i * row_h
            name = pick(phase, "name", "title", default=f"Phase {i+1}")
            duration = pick(phase, "duration", "weeks")
            deliverables = pick_list(phase, "deliverables", "outputs")
            quick_wins = pick_list(phase, "quick_wins", "milestones")

            # Alternating row background
            if i % 2 == 0:
                bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(12.0), row_h - Inches(0.08))
                bg.fill.solid()
                bg.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
                bg.line.fill.background()

            # Phase number circle accent
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.82), y + (row_h - Inches(0.38)) / 2, Inches(0.38), Inches(0.38))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            dot.text_frame.paragraphs[0].text = str(i + 1)
            dot.text_frame.paragraphs[0].font.size = Pt(11)
            dot.text_frame.paragraphs[0].font.bold = True
            dot.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            dot.text_frame.paragraphs[0].font.name = self.FONT
            dot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Phase name
            self._add_textbox(
                slide, Inches(1.35), y + Inches(0.08), label_w, Inches(0.45),
                text=name, font_size=14, bold=True, color=color,
            )

            # Duration bar
            bar_h = Inches(0.42)
            bar_y = y + (row_h - bar_h) / 2
            bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_y, bar_w * 0.55, bar_h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            bar.text_frame.paragraphs[0].text = duration
            bar.text_frame.paragraphs[0].font.size = Pt(12)
            bar.text_frame.paragraphs[0].font.bold = True
            bar.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            bar.text_frame.paragraphs[0].font.name = self.FONT
            bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Deliverables — listed to the right of the bar
            items = quick_wins + deliverables
            if items:
                items_text = "  •  ".join(items[:3])
                self._add_textbox(
                    slide, bar_left + bar_w * 0.57, bar_y - Inches(0.02),
                    bar_w * 0.43, bar_h + Inches(0.05),
                    text=items_text, font_size=11, color=Colors.SECONDARY_TEXT,
                )

    def add_recommendation_slide(self, summary: str, next_steps: list = None):
        """Slide 9: Recommendation — dark blank slide with controlled fonts."""
        slide = self._add_blank_slide(dark=True)

        # "Recommendation" label in gold
        self._add_textbox(
            slide, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.5),
            text="Our Recommendation", font_size=13, bold=True,
            color=Colors.GOLD, font_name=self.FONT,
        )

        # Gold accent bar
        acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.05))
        acc.fill.solid()
        acc.fill.fore_color.rgb = Colors.GOLD
        acc.line.fill.background()

        # Main recommendation statement — 20pt, white, readable
        rec_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(2.2))
        tf = rec_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = summary
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.name = self.FONT_HEADING
        p.font.color.rgb = Colors.WHITE

        # Next steps section
        if next_steps:
            self._add_textbox(
                slide, Inches(0.8), Inches(3.6), Inches(4.0), Inches(0.4),
                text="NEXT STEPS", font_size=11, bold=True,
                color=Colors.GOLD,
            )
            # Divider
            div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.05), Inches(11.7), Inches(0.03))
            div.fill.solid()
            div.fill.fore_color.rgb = RGBColor(0x55, 0x50, 0x4C)
            div.line.fill.background()

            y = Inches(4.15)
            for i, step in enumerate(next_steps[:4]):
                action = step.get("action", str(step)) if isinstance(step, dict) else str(step)
                owner = step.get("owner", "") if isinstance(step, dict) else ""
                deadline = step.get("deadline", "") if isinstance(step, dict) else ""

                # Number
                self._add_textbox(
                    slide, Inches(0.8), y, Inches(0.4), Inches(0.5),
                    text=str(i + 1), font_size=13, bold=True, color=Colors.GOLD,
                )
                # Action text
                action_text = action
                meta_parts = []
                if owner:
                    meta_parts.append(owner)
                if deadline:
                    meta_parts.append(deadline)
                meta = f"  [{' · '.join(meta_parts)}]" if meta_parts else ""
                self._add_textbox(
                    slide, Inches(1.25), y, Inches(11.2), Inches(0.5),
                    text=action_text + meta, font_size=13, color=Colors.WHITE,
                )
                y += Inches(0.6)

    def add_disclaimer_slide(self, disclaimer: str):
        """Final slide: commercial disclaimer for price-bearing business cases."""
        if not disclaimer:
            return
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Commercial Disclaimer")
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.7), Inches(4.9))
        tf = body.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = disclaimer
        p.font.size = Pt(12)
        p.font.name = self.FONT
        p.font.color.rgb = Colors.PRIMARY_TEXT

    # ================================================================
    # Build from YAML spec
    # ================================================================

    @classmethod
    def from_spec(cls, spec: dict, template: Optional[str] = None) -> "BusinessCaseDeckGenerator":
        """Build a complete business case deck from a YAML specification."""
        spec = enrich_adbs_to_adbd_business_case(spec)
        spec = _enrich_sparse_business_case(spec)
        bc = spec.get("business_case", spec)  # Support both wrapped and unwrapped
        gen = cls(template=template)

        # Slide 1: Cover
        gen.add_cover_slide(
            customer=pick(bc, "customer_name", "customer"),
            subtitle=pick(
                bc,
                "cover_subtitle",
                "subtitle",
                default="ADB-S to ADB-D Migration Business Case",
            ),
            prepared_by=pick(bc, "prepared_by", "author"),
            prepared_by_role=pick(bc, "prepared_by_role", "role", "title"),
            date=pick(bc, "date", "generated_on"),
        )

        # Slide 2: Executive Summary
        exec_summary = bc.get("executive_summary", "")
        if exec_summary:
            gen.add_executive_summary_slide(exec_summary)

        # Slide 3: Business Drivers
        drivers_data = bc.get("drivers", {})
        # Accept primary_driver at the top level as well (common user shape)
        top_primary = pick(bc, "primary_driver", "main_driver")
        # Accept a plain string or list for drivers
        if isinstance(drivers_data, str):
            drivers_data = {"primary": drivers_data}
        elif isinstance(drivers_data, list):
            drivers_data = {"items": drivers_data}
        elif not isinstance(drivers_data, dict):
            drivers_data = {}

        if drivers_data or top_primary:
            driver_statements = []
            primary = pick(drivers_data, "primary", "primary_driver", "main_driver") or top_primary
            urgency = pick(drivers_data, "urgency", "why_now")
            coi = drivers_data.get("cost_of_inaction", {}) or {}

            if primary:
                primary_text = primary if isinstance(primary, str) else str(primary)
                # Only rewrite snake_case enum tokens ("compliance_driven"); leave natural
                # language intact so values like "Soberanía de datos + compliance regulatorio"
                # render verbatim.
                if "_" in primary_text and " " not in primary_text and len(primary_text) <= 40:
                    headline = f"Primary Driver: {primary_text.replace('_', ' ').title()}"
                else:
                    headline = primary_text
                driver_statements.append(f"{headline}\n{urgency}" if urgency else headline)

            # Spec-provided additional cards take precedence over the inaction fallback
            extra_items = pick_list(drivers_data, "items", "additional", "secondary", "cards")
            for item in extra_items:
                if isinstance(item, str) and item.strip():
                    driver_statements.append(item.strip())
                elif isinstance(item, dict):
                    label = pick(item, "label", "headline", "title", "name")
                    desc = pick(item, "detail", "description", "text", "body")
                    if label and desc:
                        driver_statements.append(f"{label}\n{desc}")
                    elif label or desc:
                        driver_statements.append(label or desc)

            if not extra_items:
                financial = pick(coi, "financial")
                operational = pick(coi, "operational")
                strategic = pick(coi, "strategic")
                if financial:
                    driver_statements.append(f"Financial Impact of Inaction\n{financial}")
                if operational:
                    driver_statements.append(f"Operational Impact\n{operational}")
                elif strategic:
                    driver_statements.append(f"Strategic Risk\n{strategic}")

            if driver_statements:
                gen.add_business_drivers_slide(driver_statements)

        # Slide 4: TCO Comparison
        tco = bc.get("tco", {})
        if tco and (tco.get("current_state") or tco.get("proposed_oci")):
            gen.add_tco_slide(tco)

        breakdown = tco.get("breakdown", {})
        if breakdown:
            gen.add_cost_breakdown_slide(breakdown)

        projection = pick_list(tco, "projection", "forecast")
        if projection:
            gen.add_tco_projection_slide(
                projection,
                title=pick(tco, "projection_title", default="3-Year TCO Projection"),
                storage_economics=tco.get("storage_economics"),
            )

        crossover_chart = tco.get("crossover_chart", {})
        if crossover_chart:
            gen.add_tco_crossover_slide(crossover_chart)

        value_case = tco.get("business_value", {})
        if value_case:
            gen.add_business_value_slide(value_case)

        # Slide 5: ROI
        roi = bc.get("roi", {})
        if roi and (roi.get("cards") or any(roi.get(k) for k in ["three_year_roi_pct", "payback_months", "total_investment"])):
            gen.add_roi_slide(roi, headline=pick(roi, "headline"))

        # Slide 6: Value Drivers
        value_drivers = bc.get("value_drivers", [])
        if value_drivers:
            gen.add_value_drivers_slide(value_drivers)

        # Slide 7: Risk Assessment
        risks = bc.get("risks", {})
        migration_risks = pick_list(risks, "migration_risks")
        do_nothing_risks = pick_list(risks, "do_nothing_risks")
        if migration_risks or do_nothing_risks:
            gen.add_risk_slide(migration_risks, do_nothing_risks)

        # Slide 8: Roadmap
        roadmap = bc.get("roadmap", {})
        if roadmap.get("phases"):
            gen.add_roadmap_slide(
                phases=pick_list(roadmap, "phases"),
                total_duration=pick(roadmap, "total_duration"),
            )

        # Slide 9: Recommendation
        rec = bc.get("recommendation", {})
        if rec:
            summary = pick(rec, "summary")
            if not summary and rec.get("commitment_amount"):
                summary = f"Recommended: {rec['commitment_amount']} {rec.get('commitment_type', 'UCM')} commitment"
            if summary:
                gen.add_recommendation_slide(
                    summary=summary,
                    next_steps=pick_list(rec, "next_steps"),
                )

        disclaimer = pick(bc, "commercial_disclaimer", "disclaimer")
        if disclaimer:
            gen.add_disclaimer_slide(disclaimer)

        return gen


# ============================================================
# CLI
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate OCI business case slide deck (.pptx)"
    )
    parser.add_argument(
        "--spec", required=True,
        help="Path to YAML business case spec file",
    )
    parser.add_argument(
        "--output", default="business-case.pptx",
        help="Output .pptx file path",
    )
    parser.add_argument(
        "--template",
        help="Path to Oracle FY26 .pptx template (default: templates/Oracle_PPT-template_FY26.pptx)",
    )
    parser.add_argument(
        "--no-boms",
        action="store_true",
        help="Do not generate companion BOM YAML/XLSX files for ADB-S to ADB-D cases.",
    )
    parser.add_argument(
        "--no-validate",
        action="store_true",
        help="Skip structural PPTX validation.",
    )
    args = parser.parse_args()

    with open(args.spec, 'r') as f:
        spec = yaml.safe_load(f)

    bc = spec.get("business_case", spec)
    if not pick(bc, "commercial_disclaimer", "disclaimer"):
        catalog_path = Path(__file__).resolve().parent.parent / "kb" / "pricing" / "oci-sku-catalog.yaml"
        if catalog_path.exists():
            catalog = yaml.safe_load(catalog_path.read_text(encoding="utf-8")) or {}
            disclaimer = catalog.get("disclaimer", "")
            if disclaimer:
                bc["commercial_disclaimer"] = disclaimer

    gen = BusinessCaseDeckGenerator.from_spec(spec, template=args.template)
    gen.save(args.output)

    print(f"Generated: {args.output}")
    print(f"  Slides: {gen.slide_count}")
    print(f"  Template: Oracle FY26")
    customer = spec.get("business_case", spec).get("customer_name", "")
    if customer:
        print(f"  Customer: {customer}")

    adbd_config = bc.get("adbs_to_adbd") or bc.get("adb_s_to_adb_d")
    if isinstance(adbd_config, dict) and not args.no_boms:
        bom_dir = Path(args.output).with_suffix("").parent / f"{Path(args.output).stem}-boms"
        adbd_config = {**adbd_config, "customer_name": pick(bc, "customer_name", "customer"), "prepared_by": pick(bc, "prepared_by", "author")}
        saved = write_bom_outputs(adbd_config, bom_dir)
        print(f"  BOM outputs: {len(saved)} files in {bom_dir}")

    if not args.no_validate:
        validation_cfg = bc.get("deck_validation", {}) if isinstance(bc.get("deck_validation"), dict) else {}
        required = validation_cfg.get("required_phrases")
        forbidden = validation_cfg.get("forbidden_phrases")
        expected_titles = validation_cfg.get("expected_titles")
        if required is None and isinstance(adbd_config, dict):
            required = [
                "BYOL/PAYG model",
                "Discount",
                "GoldenGate bridge duration",
                "Workload ECPU demand",
                "ECPU capacity",
                "Storage break-even",
                "Crossover",
            ]
        if forbidden is None:
            forbidden = ["OCI Annual", "FTE-year"]
        if expected_titles is None and isinstance(adbd_config, dict):
            crossover_title = pick(adbd_config.get("crossover_chart") or {}, "title", default="TCO Crossover")
            expected_titles = [
                "Total Cost of Ownership",
                "BOM + Operations Cost Breakdown",
                crossover_title,
                "Business Value Model",
                "Commercial Disclaimer",
            ]
        report = validate_pptx(
            args.output,
            expected_titles=expected_titles or [],
            required_phrases=required or [],
            forbidden_phrases=forbidden or [],
            min_slides=validation_cfg.get("min_slides"),
            max_slides=validation_cfg.get("max_slides"),
            exact_slides=validation_cfg.get("exact_slides"),
            disclaimer_last=bool(bc.get("commercial_disclaimer")),
        )
        print(f"  Deck validation: {report['status']} ({report['slide_count']} slides)")
        for issue in report["issues"]:
            print(f"    - {issue['code']}: {issue['message']}")


if __name__ == "__main__":
    main()
