#!/usr/bin/env python3
"""Structural PPTX text validator.

Reads text boxes, table cells, and grouped shapes so lightweight preview
limitations do not hide real deck content from validation.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _shape_text(shape) -> list[str]:
    texts: list[str] = []
    if hasattr(shape, "text") and str(shape.text or "").strip():
        texts.append(str(shape.text).strip())
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                value = str(cell.text or "").strip()
                if value:
                    texts.append(value)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            texts.extend(_shape_text(child))
    elif hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(_shape_text(child))
    return texts


def extract_pptx_text(path: str | Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    slides = []
    all_text: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(_shape_text(shape))
        clean = [text for text in texts if text]
        slides.append({"index": idx, "text": "\n".join(clean), "texts": clean})
        all_text.extend(clean)
    return {"slide_count": len(prs.slides), "slides": slides, "text": "\n".join(all_text)}


def _contains(haystack: str, needle: str, case_sensitive: bool = False) -> bool:
    if case_sensitive:
        return needle in haystack
    return needle.lower() in haystack.lower()


def validate_pptx(
    path: str | Path,
    *,
    expected_titles: list[str] | None = None,
    required_phrases: list[str] | None = None,
    forbidden_phrases: list[str] | None = None,
    min_slides: int | None = None,
    max_slides: int | None = None,
    exact_slides: int | None = None,
    disclaimer_last: bool = False,
    case_sensitive: bool = False,
) -> dict[str, Any]:
    extracted = extract_pptx_text(path)
    text = extracted["text"]
    issues: list[dict[str, str]] = []

    slide_count = extracted["slide_count"]
    if exact_slides is not None and slide_count != exact_slides:
        issues.append({"code": "SLIDE_COUNT", "message": f"Expected {exact_slides} slides, found {slide_count}."})
    if min_slides is not None and slide_count < min_slides:
        issues.append({"code": "SLIDE_COUNT_MIN", "message": f"Expected at least {min_slides} slides, found {slide_count}."})
    if max_slides is not None and slide_count > max_slides:
        issues.append({"code": "SLIDE_COUNT_MAX", "message": f"Expected at most {max_slides} slides, found {slide_count}."})

    for title in expected_titles or []:
        if not _contains(text, title, case_sensitive):
            issues.append({"code": "MISSING_TITLE", "message": f"Expected title not found: {title}"})

    for phrase in required_phrases or []:
        if not _contains(text, phrase, case_sensitive):
            issues.append({"code": "MISSING_REQUIRED", "message": f"Required phrase not found: {phrase}"})

    for phrase in forbidden_phrases or []:
        if _contains(text, phrase, case_sensitive):
            issues.append({"code": "FORBIDDEN_PHRASE", "message": f"Forbidden phrase found: {phrase}"})

    if disclaimer_last:
        if not extracted["slides"]:
            issues.append({"code": "DISCLAIMER_LAST", "message": "Deck has no slides."})
        else:
            last_text = extracted["slides"][-1]["text"]
            if not _contains(last_text, "Commercial Disclaimer", case_sensitive):
                issues.append({"code": "DISCLAIMER_LAST", "message": "Commercial Disclaimer is not the last slide."})

    return {
        "path": str(path),
        "status": "fail" if issues else "pass",
        "slide_count": slide_count,
        "issues": issues,
        "titles_checked": expected_titles or [],
        "required_checked": required_phrases or [],
        "forbidden_checked": forbidden_phrases or [],
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Validate PPTX structure and text content.")
    parser.add_argument("--pptx", required=True, help="Path to PPTX file.")
    parser.add_argument("--expected-title", action="append", default=[])
    parser.add_argument("--required", action="append", default=[])
    parser.add_argument("--forbidden", action="append", default=[])
    parser.add_argument("--min-slides", type=int)
    parser.add_argument("--max-slides", type=int)
    parser.add_argument("--exact-slides", type=int)
    parser.add_argument("--disclaimer-last", action="store_true")
    parser.add_argument("--json", action="store_true")
    args = parser.parse_args()

    report = validate_pptx(
        args.pptx,
        expected_titles=args.expected_title,
        required_phrases=args.required,
        forbidden_phrases=args.forbidden,
        min_slides=args.min_slides,
        max_slides=args.max_slides,
        exact_slides=args.exact_slides,
        disclaimer_last=args.disclaimer_last,
    )
    if args.json:
        print(json.dumps(report, indent=2))
    else:
        print(f"pptx validation {report['status']}: {report['slide_count']} slides")
        for issue in report["issues"]:
            print(f"  - {issue['code']}: {issue['message']}")
    return 1 if report["status"] == "fail" else 0


if __name__ == "__main__":
    raise SystemExit(main())
