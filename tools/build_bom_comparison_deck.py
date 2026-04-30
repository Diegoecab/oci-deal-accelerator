#!/usr/bin/env python3
"""
Build a short PowerPoint deck comparing multiple OCI BOM scenarios.

The input spec references one or more BOM YAML specs. The script resolves
monthly and 12-month totals from the live SKU catalog and produces a
customer-ready Oracle FY26 deck with native PowerPoint charts.
"""

import argparse
import sys
from dataclasses import dataclass
from decimal import Decimal, ROUND_HALF_UP
from pathlib import Path

import yaml
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

try:
    from oci_pptx_base import Colors, Layouts, OraclePresBase
except ModuleNotFoundError:
    from tools.oci_pptx_base import Colors, Layouts, OraclePresBase


PROJECT_ROOT = Path(__file__).resolve().parent.parent
CATALOG_PATH = PROJECT_ROOT / "kb" / "pricing" / "oci-sku-catalog.yaml"
DB_SERVER_SKU = "B110627"
STORAGE_SERVER_SKU = "B110629"


def _round_money(value: Decimal) -> Decimal:
    return value.quantize(Decimal("0.01"), rounding=ROUND_HALF_UP)


def _money_text(value: Decimal) -> str:
    return f"USD {value:,.2f}"


def _pct_text(value: Decimal) -> str:
    sign = "+" if value > 0 else ""
    return f"{sign}{value:.1f}%"


def _pct_abs_text(value: Decimal) -> str:
    return f"{abs(value):.1f}%"


def _resolve_path(path_str: str) -> Path:
    path = Path(path_str)
    if path.is_absolute():
        return path
    return PROJECT_ROOT / path


def _load_catalog(path: Path) -> dict[str, dict]:
    raw = yaml.safe_load(path.read_text(encoding="utf-8"))
    catalog: dict[str, dict] = {}
    for category in raw.get("categories", {}).values():
        for sku in category.get("skus", []):
            catalog[str(sku["sku"])] = sku
    return catalog


@dataclass
class ScenarioMetrics:
    scenario_id: str
    short_label: str
    title: str
    topology: str
    tradeoff: str
    accent: str
    term_months: int
    db_servers: int
    storage_servers: int
    ecpus: int
    monthly_total: Decimal
    term_total: Decimal
    infra_monthly: Decimal
    compute_monthly: Decimal
    delta_vs_baseline: Decimal = Decimal("0")
    delta_pct_vs_baseline: Decimal = Decimal("0")


def _scenario_color(accent: str) -> RGBColor:
    return {
        "recommended": Colors.COPPER,
        "balanced": Colors.TEAL,
        "savings": Colors.FOREST,
    }.get(accent, Colors.TEAL)


def _load_bom_metrics(spec_path: Path, meta: dict, catalog: dict[str, dict]) -> ScenarioMetrics:
    raw = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    bom = raw.get("bom", raw)
    items = bom.get("line_items", [])

    monthly_total = Decimal("0")
    term_total = Decimal("0")
    infra_monthly = Decimal("0")
    compute_monthly = Decimal("0")
    db_servers = Decimal("0")
    storage_servers = Decimal("0")
    ecpus = Decimal("0")
    months_seen: set[int] = set()

    for item in items:
        sku = str(item["sku"])
        entry = catalog.get(sku)
        if not entry:
            raise KeyError(f"SKU {sku} not found in catalog while processing {spec_path}")

        price = Decimal(str(entry.get("list_price_usd", 0)))
        qty = Decimal(str(item.get("qty", 0)))
        hours_units = Decimal(str(item.get("hours_units", entry.get("default_hours_units", 730))))
        months = int(item.get("months", 12))
        discount = Decimal(str(item.get("discount", 0)))
        line_monthly = price * qty * hours_units * (Decimal("1") - discount)
        line_term = line_monthly * Decimal(str(months))

        monthly_total += line_monthly
        term_total += line_term
        months_seen.add(months)

        if sku == DB_SERVER_SKU:
            db_servers += qty
            infra_monthly += line_monthly
        elif sku == STORAGE_SERVER_SKU:
            storage_servers += qty
            infra_monthly += line_monthly
        elif "ECPU" in str(entry.get("metric", "")).upper():
            ecpus += qty
            compute_monthly += line_monthly
        else:
            compute_monthly += line_monthly

    term_months = max(months_seen) if months_seen else int(meta.get("term_months", 12))
    return ScenarioMetrics(
        scenario_id=meta["id"],
        short_label=meta.get("short_label", meta["id"]),
        title=meta["title"],
        topology=meta["topology"],
        tradeoff=meta["tradeoff"],
        accent=meta.get("accent", "balanced"),
        term_months=term_months,
        db_servers=int(db_servers),
        storage_servers=int(storage_servers),
        ecpus=int(ecpus),
        monthly_total=_round_money(monthly_total),
        term_total=_round_money(term_total),
        infra_monthly=_round_money(infra_monthly),
        compute_monthly=_round_money(compute_monthly),
    )


class BOMComparisonDeck(OraclePresBase):
    TITLE_ACCENT_COLOR = Colors.TEAL
    MARGIN = Inches(0.55)

    def __init__(self, spec: dict, scenarios: list[ScenarioMetrics]):
        super().__init__()
        self.spec = spec
        self.deck_meta = spec["deck"]
        self.scenarios = scenarios
        self.term_months = int(self.deck_meta.get("term_months", 12))

    @property
    def _compute_is_constant(self) -> bool:
        return len({scenario.compute_monthly for scenario in self.scenarios}) == 1

    def _add_footer_note(self, slide, text: str) -> None:
        box = self._add_textbox(
            slide,
            self.MARGIN,
            Inches(7.0),
            Inches(12.0),
            Inches(0.25),
            text=text,
            font_size=9,
            italic=True,
            color=Colors.SECONDARY_TEXT,
        )
        box.text_frame.word_wrap = True

    def _add_pill(self, slide, left, top, width, text, fill_color: RGBColor) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.FONT
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = Colors.WHITE
        shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def add_cover(self) -> None:
        slide = self._add_layout_slide(Layouts.COVER_DARK)
        self._set_placeholder(slide, 0, self.deck_meta["title"])
        self._set_placeholder(slide, 1, self.deck_meta["subtitle"])
        disclaimer = self.deck_meta.get("commercial_disclaimer", {})
        cover_tag = disclaimer.get("cover_tag")
        if cover_tag:
            self._add_pill(slide, Inches(8.7), Inches(0.45), Inches(3.9), cover_tag, Colors.COPPER)

    def add_summary_slide(self) -> None:
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Scenario Summary", margin=self.MARGIN)

        if self._compute_is_constant:
            headline = (
                "All three scenarios carry the same 830 ECPU allocation. "
                "The cost delta comes from infrastructure footprint only."
            )
        else:
            headline = (
                "Recommended setup remains the architecture-led option. "
                "Combined 6x4 trims infrastructure. Combined 5x3 is the strongest cost-down option."
            )
        self._add_textbox(
            slide,
            self.MARGIN,
            Inches(1.2),
            Inches(12.0),
            Inches(0.5),
            text=headline,
            font_size=15,
            bold=True,
            color=Colors.PRIMARY_TEXT,
        )

        card_lefts = [self.MARGIN, Inches(4.45), Inches(8.35)]
        card_width = Inches(3.85)
        card_top = Inches(1.9)
        card_height = Inches(4.55)
        baseline = next(s for s in self.scenarios if s.scenario_id == self.deck_meta["baseline_scenario_id"])

        for left, scenario in zip(card_lefts, self.scenarios):
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = Colors.WARM_BG
            card.line.color.rgb = _scenario_color(scenario.accent)
            card.line.width = Pt(1.5)

            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, card_width, Inches(0.18))
            accent.fill.solid()
            accent.fill.fore_color.rgb = _scenario_color(scenario.accent)
            accent.line.fill.background()

            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(0.28), card_width - Inches(0.36), Inches(0.55),
                text=scenario.title, font_size=18, bold=True, color=Colors.PRIMARY_TEXT, font_name=self.FONT_HEADING,
            )
            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(0.82), card_width - Inches(0.36), Inches(0.7),
                text=scenario.topology, font_size=11, color=Colors.SECONDARY_TEXT,
            )

            tag_text = "Recommended" if scenario.scenario_id == baseline.scenario_id else (
                "Lowest cost" if scenario.monthly_total == min(s.monthly_total for s in self.scenarios) else "Alternative"
            )
            self._add_pill(slide, left + Inches(0.18), card_top + Inches(1.45), Inches(1.35), tag_text, _scenario_color(scenario.accent))

            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(1.92), card_width - Inches(0.36), Inches(0.32),
                text="Monthly", font_size=10, bold=True, color=Colors.SECONDARY_TEXT,
            )
            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(2.18), card_width - Inches(0.36), Inches(0.42),
                text=_money_text(scenario.monthly_total), font_size=20, bold=True, color=_scenario_color(scenario.accent),
            )
            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(2.65), card_width - Inches(0.36), Inches(0.3),
                text=f"{self.term_months}M total: {_money_text(scenario.term_total)}", font_size=10, color=Colors.SECONDARY_TEXT,
            )

            if scenario.scenario_id == baseline.scenario_id:
                delta_text = "Baseline"
                delta_color = Colors.PRIMARY_TEXT
            elif scenario.delta_vs_baseline < 0:
                delta_text = (
                    f"Savings vs {baseline.short_label}: "
                    f"{_money_text(-scenario.delta_vs_baseline)} / mo "
                    f"({_pct_abs_text(scenario.delta_pct_vs_baseline)} lower)"
                )
                delta_color = Colors.FOREST
            else:
                delta_text = (
                    f"Premium vs {baseline.short_label}: "
                    f"{_money_text(scenario.delta_vs_baseline)} / mo "
                    f"({_pct_text(scenario.delta_pct_vs_baseline)})"
                )
                delta_color = Colors.PRIMARY_TEXT
            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(3.05), card_width - Inches(0.36), Inches(0.45),
                text=delta_text, font_size=10, bold=True,
                color=delta_color,
            )

            footprint = (
                f"DB servers: {scenario.db_servers}\n"
                f"Storage servers: {scenario.storage_servers}\n"
                f"Total ECPU: {scenario.ecpus:,}"
            )
            tf_box = self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(3.55), card_width - Inches(0.36), Inches(0.72),
                text=footprint, font_size=10, color=Colors.PRIMARY_TEXT,
            )
            tf_box.text_frame.word_wrap = True

            self._add_textbox(
                slide, left + Inches(0.18), card_top + Inches(4.15), card_width - Inches(0.36), Inches(0.32),
                text=scenario.tradeoff, font_size=10, italic=True, color=Colors.SECONDARY_TEXT,
            )

        self._add_footer_note(slide, "Source footprint: recommended and combined options from the reference deck. Pricing based on OCI list price, 730 hours/month, no commercial discount.")

    def add_monthly_chart_slide(self) -> None:
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Monthly Cost Comparison", margin=self.MARGIN)

        data = CategoryChartData()
        data.categories = [s.short_label for s in self.scenarios]
        data.add_series("Monthly USD", [float(s.monthly_total) for s in self.scenarios])
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            self.MARGIN,
            Inches(1.35),
            Inches(7.55),
            Inches(4.9),
            data,
        )
        chart = frame.chart
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.number_format = '$0,"K"'
        chart.value_axis.format.line.color.rgb = Colors.SECONDARY_TEXT
        chart.category_axis.format.line.color.rgb = Colors.SECONDARY_TEXT

        plot = chart.plots[0]
        plot.vary_by_categories = True
        plot.has_data_labels = True
        plot.data_labels.number_format = '$#,##0'
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        series = chart.series[0]
        for point, scenario in zip(series.points, self.scenarios):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _scenario_color(scenario.accent)
            point.format.line.color.rgb = _scenario_color(scenario.accent)

        callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.45), Inches(4.5), Inches(4.6))
        callout.fill.solid()
        callout.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
        callout.line.color.rgb = Colors.TEAL
        callout.line.width = Pt(1.2)

        baseline = next(s for s in self.scenarios if s.scenario_id == self.deck_meta["baseline_scenario_id"])
        cheapest = min(self.scenarios, key=lambda s: s.monthly_total)
        bullets = [
            f"{baseline.short_label} is still the architecture-led recommendation.",
            f"{self.scenarios[1].short_label} saves {_money_text(-self.scenarios[1].delta_vs_baseline)} per month by reducing storage footprint only.",
            f"{cheapest.short_label} saves {_money_text(-cheapest.delta_vs_baseline)} per month and is {_pct_abs_text(cheapest.delta_pct_vs_baseline)} below the baseline.",
            (
                f"{cheapest.short_label} is cheaper because it cuts both infrastructure and compute footprint."
                if not self._compute_is_constant
                else f"{cheapest.short_label} is cheaper exclusively because it uses the smallest infrastructure footprint."
            ),
        ]
        self._add_textbox(
            slide, Inches(8.42), Inches(1.7), Inches(3.95), Inches(0.35),
            text="Key takeaways", font_size=16, bold=True, color=Colors.PRIMARY_TEXT, font_name=self.FONT_HEADING,
        )
        notes_box = self._add_textbox(slide, Inches(8.35), Inches(2.15), Inches(3.95), Inches(3.4), text="", font_size=11)
        tf = notes_box.text_frame
        tf.clear()
        for idx, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = self.FONT
            p.font.size = Pt(11)
            p.font.color.rgb = Colors.PRIMARY_TEXT
            p.space_after = Pt(8)

        self._add_footer_note(slide, "This chart shows monthly run-rate only. The same ranking holds for the 12-month view because all scenarios use the same 730-hour monthly convention and 12-month term.")

    def add_composition_slide(self) -> None:
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Cost Composition: Infrastructure vs ECPU", margin=self.MARGIN)

        data = CategoryChartData()
        data.categories = [s.short_label for s in self.scenarios]
        data.add_series("Infrastructure", [float(s.infra_monthly) for s in self.scenarios])
        data.add_series("ADB Dedicated ECPU BYOL", [float(s.compute_monthly) for s in self.scenarios])

        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            self.MARGIN,
            Inches(1.35),
            Inches(7.55),
            Inches(4.8),
            data,
        )
        chart = frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.number_format = '$0,"K"'

        series_infra = chart.series[0]
        series_infra.format.fill.solid()
        series_infra.format.fill.fore_color.rgb = Colors.MUTED_TEAL
        series_infra.format.line.color.rgb = Colors.MUTED_TEAL

        series_compute = chart.series[1]
        series_compute.format.fill.solid()
        series_compute.format.fill.fore_color.rgb = Colors.COPPER
        series_compute.format.line.color.rgb = Colors.COPPER

        table = self._add_table(
            slide,
            len(self.scenarios) + 1,
            4,
            Inches(8.2),
            Inches(1.6),
            Inches(4.4),
            Inches(2.4),
        )
        headers = ["Scenario", "Infra / mo", "ECPU / mo", "Infra share"]
        widths = [1.2, 1.05, 1.1, 1.0]
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)
        for col, header in enumerate(headers):
            self._style_table_cell(
                table.cell(0, col), header,
                font_size=10, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT,
            )

        for row_idx, scenario in enumerate(self.scenarios, start=1):
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            infra_share = (scenario.infra_monthly / scenario.monthly_total * Decimal("100")) if scenario.monthly_total else Decimal("0")
            self._style_table_cell(table.cell(row_idx, 0), scenario.short_label, font_size=10, bold=True, bg_color=bg)
            self._style_table_cell(table.cell(row_idx, 1), _money_text(scenario.infra_monthly), font_size=9, bg_color=bg, alignment=PP_ALIGN.RIGHT)
            self._style_table_cell(table.cell(row_idx, 2), _money_text(scenario.compute_monthly), font_size=9, bg_color=bg, alignment=PP_ALIGN.RIGHT)
            self._style_table_cell(table.cell(row_idx, 3), f"{infra_share:.1f}%", font_size=9, bg_color=bg, alignment=PP_ALIGN.RIGHT)

        insight = (
            "The ECPU line is flat across all three scenarios: 830 total ECPU.\n"
            "Only the DB server and storage server counts move the price."
            if self._compute_is_constant
            else
            "Combined 6x4 changes only the infrastructure layer.\n"
            "Combined 5x3 is materially cheaper because it also reduces the ECPU footprint."
        )
        callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(4.35), Inches(4.4), Inches(1.15))
        callout.fill.solid()
        callout.fill.fore_color.rgb = Colors.WARM_BG
        callout.line.color.rgb = Colors.COPPER
        self._add_textbox(
            slide, Inches(8.42), Inches(4.58), Inches(4.0), Inches(0.7),
            text=insight, font_size=11, bold=True, color=Colors.PRIMARY_TEXT,
        )

        self._add_footer_note(slide, "Infrastructure is the sum of X11M DB servers plus X11M storage servers. Compute is the ADB Dedicated BYOL ECPU line.")

    def add_detail_table_slide(self) -> None:
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Detailed BOM Comparison", margin=self.MARGIN)

        table = self._add_table(
            slide,
            len(self.scenarios) + 1,
            8,
            self.MARGIN,
            Inches(1.4),
            Inches(12.2),
            Inches(3.1),
        )
        widths = [1.5, 3.0, 0.9, 1.1, 1.1, 1.35, 1.35, 1.9]
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)
        headers = ["Scenario", "Topology", "DB", "Storage", "ECPU", "Monthly", f"{self.term_months}M total", "Trade-off"]
        for col, header in enumerate(headers):
            self._style_table_cell(
                table.cell(0, col), header,
                font_size=10, bold=True, color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if col not in (0, 1, 7) else PP_ALIGN.LEFT,
            )

        for row_idx, scenario in enumerate(self.scenarios, start=1):
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            if scenario.scenario_id == self.deck_meta["baseline_scenario_id"]:
                bg = Colors.TEAL
                txt = Colors.WHITE
            else:
                txt = Colors.PRIMARY_TEXT

            values = [
                scenario.short_label,
                scenario.topology,
                str(scenario.db_servers),
                str(scenario.storage_servers),
                f"{scenario.ecpus:,}",
                _money_text(scenario.monthly_total),
                _money_text(scenario.term_total),
                scenario.tradeoff,
            ]
            aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT]
            for col_idx, (value, align) in enumerate(zip(values, aligns)):
                self._style_table_cell(
                    table.cell(row_idx, col_idx), value,
                    font_size=9, bold=(scenario.scenario_id == self.deck_meta["baseline_scenario_id"]),
                    color=txt, bg_color=bg, alignment=align,
                )

        deltas_box = self._add_textbox(slide, self.MARGIN, Inches(5.0), Inches(12.0), Inches(1.0), text="", font_size=10)
        tf = deltas_box.text_frame
        tf.clear()
        baseline = next(s for s in self.scenarios if s.scenario_id == self.deck_meta["baseline_scenario_id"])
        summary_lines = [
            f"Baseline: {baseline.short_label} at {_money_text(baseline.monthly_total)} per month.",
            f"{self.scenarios[1].short_label}: {_money_text(-self.scenarios[1].delta_vs_baseline)} per month lower than baseline.",
            f"{self.scenarios[2].short_label}: {_money_text(-self.scenarios[2].delta_vs_baseline)} per month lower than baseline.",
        ]
        for idx, line in enumerate(summary_lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = self.FONT
            p.font.size = Pt(10)
            p.font.color.rgb = Colors.SECONDARY_TEXT
            p.space_after = Pt(5)

        self._add_footer_note(slide, "BOM workbook references are saved alongside this deck for traceability and recalculation.")

    def add_disclaimer_slide(self) -> None:
        disclaimer = self.deck_meta.get("commercial_disclaimer", {})
        if not disclaimer:
            return

        slide = self._add_blank_slide()
        self._add_title_bar(slide, disclaimer.get("title", "Commercial Disclaimer"), margin=self.MARGIN)

        intro_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.MARGIN,
            Inches(1.35),
            Inches(12.1),
            Inches(2.15),
        )
        intro_box.fill.solid()
        intro_box.fill.fore_color.rgb = Colors.WARM_BG
        intro_box.line.color.rgb = Colors.COPPER
        intro_box.line.width = Pt(1.2)

        formal = disclaimer.get("formal_text", "")
        formal_tb = self._add_textbox(
            slide,
            Inches(0.8),
            Inches(1.58),
            Inches(11.55),
            Inches(1.7),
            text="",
            font_size=11,
            color=Colors.PRIMARY_TEXT,
        )
        formal_tf = formal_tb.text_frame
        formal_tf.clear()
        p = formal_tf.paragraphs[0]
        p.text = formal
        p.font.name = self.FONT
        p.font.size = Pt(11)
        p.font.color.rgb = Colors.PRIMARY_TEXT
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(0)
        formal_tf.word_wrap = True

        note_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.MARGIN,
            Inches(3.95),
            Inches(12.1),
            Inches(2.15),
        )
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = Colors.TABLE_ALT_ROW
        note_box.line.color.rgb = Colors.TEAL
        note_box.line.width = Pt(1.2)

        self._add_textbox(
            slide,
            Inches(0.8),
            Inches(4.18),
            Inches(3.2),
            Inches(0.35),
            text="Practical notes",
            font_size=16,
            bold=True,
            color=Colors.PRIMARY_TEXT,
            font_name=self.FONT_HEADING,
        )
        bullets_tb = self._add_textbox(
            slide,
            Inches(0.8),
            Inches(4.55),
            Inches(11.2),
            Inches(1.25),
            text="",
            font_size=11,
            color=Colors.PRIMARY_TEXT,
        )
        bullets_tf = bullets_tb.text_frame
        bullets_tf.clear()
        for idx, bullet in enumerate(disclaimer.get("bullets", [])):
            p = bullets_tf.paragraphs[0] if idx == 0 else bullets_tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.name = self.FONT
            p.font.size = Pt(11)
            p.font.color.rgb = Colors.PRIMARY_TEXT
            p.space_after = Pt(10)

        self._add_footer_note(
            slide,
            "Reference pricing only. Formal commercials must move through the Oracle quoting process and the applicable agreement path.",
        )

    def build(self) -> None:
        self.add_cover()
        self.add_summary_slide()
        self.add_monthly_chart_slide()
        self.add_composition_slide()
        self.add_detail_table_slide()
        self.add_disclaimer_slide()


def load_deck_inputs(spec_path: Path) -> tuple[dict, list[ScenarioMetrics]]:
    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    deck = spec["deck"]
    catalog = _load_catalog(CATALOG_PATH)
    scenarios = []
    for scenario_meta in deck.get("scenarios", []):
        bom_spec = _resolve_path(scenario_meta["bom_spec"])
        scenarios.append(_load_bom_metrics(bom_spec, scenario_meta, catalog))

    baseline_id = deck["baseline_scenario_id"]
    baseline = next(s for s in scenarios if s.scenario_id == baseline_id)
    for scenario in scenarios:
        scenario.delta_vs_baseline = _round_money(scenario.monthly_total - baseline.monthly_total)
        if baseline.monthly_total:
            pct = (scenario.monthly_total - baseline.monthly_total) / baseline.monthly_total * Decimal("100")
            scenario.delta_pct_vs_baseline = Decimal(str(pct)).quantize(Decimal("0.1"), rounding=ROUND_HALF_UP)
    return spec, scenarios


def main() -> int:
    parser = argparse.ArgumentParser(description="Build a comparison deck from multiple BOM specs.")
    parser.add_argument("--spec", required=True, help="YAML deck spec path")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    spec_path = _resolve_path(args.spec)
    output_path = _resolve_path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    spec, scenarios = load_deck_inputs(spec_path)
    deck = BOMComparisonDeck(spec, scenarios)
    deck.build()
    deck.save(str(output_path))
    print(f"Deck saved: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
