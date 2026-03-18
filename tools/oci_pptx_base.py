#!/usr/bin/env python3
"""
OCI Deal Accelerator — Shared PPTX base (Oracle FY26 template + Redwood design system).

Import: from oci_pptx_base import Colors, Layouts, OraclePresBase
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# Oracle Redwood Design System Colors
# Source: Oracle Database Icon Library [July 2024] — theme RMIL_Master
# Color scheme: "Oracle Redwood"
# ============================================================

class Colors:
    # Primary text — never pure black
    DARK_BG = RGBColor(0x31, 0x2D, 0x2A)       # #312D2A — dk1
    PRIMARY_TEXT = RGBColor(0x31, 0x2D, 0x2A)    # #312D2A — dk1
    SECONDARY_TEXT = RGBColor(0x6F, 0x69, 0x64)  # #6F6964 — Neutral 60
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    WARM_BG = RGBColor(0xFC, 0xFB, 0xFA)         # #FCFBFA — lt1, content bg

    # Oracle Redwood accent palette
    ORACLE_RED = RGBColor(0xC7, 0x46, 0x34)      # #C74634 — accent1, primary brand
    GOLD = RGBColor(0xFA, 0xCD, 0x62)            # #FACD62 — accent2
    TEAL = RGBColor(0x2C, 0x59, 0x67)            # #2C5967 — hlink, tables/links
    SAGE = RGBColor(0x75, 0x9C, 0x6C)            # #759C6C — accent6
    FOREST = RGBColor(0x2B, 0x62, 0x42)          # #2B6242 — accent4
    MUTED_TEAL = RGBColor(0x94, 0xAF, 0xAF)      # #94AFAF — accent3
    BURNT_ORANGE = RGBColor(0xAE, 0x56, 0x2C)    # #AE562C — accent5

    # OCI service category colors (used in architecture diagrams)
    COPPER = RGBColor(0xAA, 0x64, 0x3B)          # #AA643B — database
    PURPLE = RGBColor(0x80, 0x49, 0x98)          # #804998 — integration

    # Status
    SUCCESS = RGBColor(0x2B, 0x62, 0x42)         # #2B6242 — forest green
    WARNING = RGBColor(0xAE, 0x56, 0x2C)         # #AE562C — burnt orange
    ERROR = RGBColor(0xC7, 0x46, 0x34)           # #C74634 — Oracle red

    # Table
    TABLE_ALT_ROW = RGBColor(0xF5, 0xF4, 0xF2)  # #F5F4F2 — Neutral 10
    TABLE_HEADER_BG = RGBColor(0x2C, 0x59, 0x67) # teal


# ============================================================
# Oracle FY26 Template Layout Indices
# ============================================================

class Layouts:
    """Layout indices from Oracle_PPT-template_FY26.pptx (104 layouts)."""
    COVER_DARK = 0              # Dark - Title_Pillar
    COVER_LIGHT = 1             # Light - Title_Pillar
    DIVIDER_LIGHT = 16          # Light - Divider
    DIVIDER_DARK = 17           # Dark - Divider
    IMPACT_DARK = 26            # Dark - Impact Statement
    IMPACT_LIGHT = 27           # Light - Impact Statement
    METRIC_DARK = 32            # Dark - Single metric
    METRIC_LIGHT = 33           # Light - Single metric
    MULTI_STATEMENT_DARK = 47   # Multi Statement – Dark
    MULTI_STATEMENT_LIGHT = 48  # Multi Statement – Light
    TWO_COL_LIGHT = 66          # Light - Title/Subtitle 2 Column
    TWO_COL_DARK = 67           # Dark - Title/Subtitle 2 Column
    FOUR_ICONS_LIGHT = 82       # Light - 4 Icons, Subhead and text
    BLANK_DARK = 84             # Dark - Blank
    BLANK_LIGHT = 85            # Light - Blank
    SAFE_HARBOR_LIGHT = 96      # Light - Safe harbor - short (if exists)
    THANK_YOU_DARK = 34         # Dark - Thank You
    THANK_YOU_LIGHT = 35        # Light - Thank You
    CLOSING_LIGHT = 101         # Light - Closing Slide (Oracle logo only)
    CLOSING_DARK = 102          # Dark - Closing Slide (Oracle logo only)


# ============================================================
# Base Presentation Class
# ============================================================

class OraclePresBase:
    """Shared PPTX infrastructure for OCI Deal Accelerator generators."""

    TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "Oracle_PPT-template_FY26.pptx"
    FONT = "Oracle Sans"
    FONT_HEADING = "Georgia"
    TITLE_ACCENT_COLOR = Colors.TEAL  # subclasses can override as class attribute

    def __init__(self, template: Optional[str] = None):
        tmpl = template or str(self.TEMPLATE_PATH)
        if not Path(tmpl).is_file():
            raise FileNotFoundError(f"Template not found: {tmpl}")
        self.prs = Presentation(tmpl)
        self._clear_slides()

    def _clear_slides(self):
        """Remove all existing slides from the template (keep layouts/masters)."""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            )
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])
        # Clean sections
        from lxml import etree  # noqa: F401
        ns_p14 = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
        for sectionLst in self.prs.part._element.findall(f'.//{{{ns_p14}}}sectionLst'):
            sectionLst.getparent().remove(sectionLst)

    def _add_layout_slide(self, layout_index: int):
        """Add a slide using a specific template layout."""
        layout = self.prs.slide_layouts[layout_index]
        return self.prs.slides.add_slide(layout)

    def _add_blank_slide(self, dark: bool = False):
        """Add a blank slide (light or dark)."""
        idx = Layouts.BLANK_DARK if dark else Layouts.BLANK_LIGHT
        return self._add_layout_slide(idx)

    def _set_placeholder(self, slide, idx: int, text: str):
        """Set text in a placeholder by index. Silently skips if not found."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                ph.text = text
                return ph
        return None

    def _style_placeholder(self, slide, idx: int, text: str,
                           font_size: int = None, bold: bool = None,
                           color: RGBColor = None):
        """Set text and style a placeholder."""
        ph = self._set_placeholder(slide, idx, text)
        if ph and ph.text_frame.paragraphs:
            for p in ph.text_frame.paragraphs:
                if font_size:
                    p.font.size = Pt(font_size)
                if bold is not None:
                    p.font.bold = bold
                if color:
                    p.font.color.rgb = color
        return ph

    def _add_textbox(self, slide, left, top, width, height,
                     text="", font_size=12, bold=False, italic=False,
                     color=None, alignment=PP_ALIGN.LEFT, font_name=None):
        """Add a text box to a slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font_name or self.FONT
        p.font.color.rgb = color or Colors.PRIMARY_TEXT
        p.alignment = alignment
        return txBox

    def _add_paragraph(self, text_frame, text="", font_size=12,
                       bold=False, italic=False, color=None,
                       alignment=PP_ALIGN.LEFT, space_before=0,
                       space_after=0, bullet=False):
        """Add a paragraph to an existing text frame."""
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = self.FONT
        p.font.color.rgb = color or Colors.PRIMARY_TEXT
        p.alignment = alignment
        if space_before:
            p.space_before = Pt(space_before)
        if space_after:
            p.space_after = Pt(space_after)
        if bullet:
            p.level = 0
        return p

    def _add_rect(self, slide, left, top, width, height,
                  fill_color=None, border_color=None):
        """Add a rectangle shape."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if border_color:
            shape.line.color.rgb = border_color
        else:
            shape.line.fill.background()
        return shape

    def _add_table(self, slide, rows, cols, left, top, width, height):
        """Add a table to a slide."""
        return slide.shapes.add_table(rows, cols, left, top, width, height).table

    def _style_table_cell(self, cell, text, font_size=10, bold=False,
                          color=None, bg_color=None, alignment=PP_ALIGN.LEFT):
        """Style a table cell."""
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
            paragraph.font.name = self.FONT
            paragraph.font.color.rgb = color or Colors.PRIMARY_TEXT
            paragraph.alignment = alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    def _add_title_bar(self, slide, title_text, margin=None):
        """Add a consistent title bar at top of content slides.

        Uses self.TITLE_ACCENT_COLOR for the accent line.
        Title is 24pt bold Georgia. margin defaults to Inches(0.8).
        """
        if margin is None:
            margin = Inches(0.8)
        self._add_textbox(
            slide, margin, Inches(0.4), Inches(11.7), Inches(0.6),
            text=title_text, font_size=24, bold=True,
            font_name=self.FONT_HEADING,
        )
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, margin, Inches(1.05), Inches(11.7), Inches(0.05)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.TITLE_ACCENT_COLOR
        bar.line.fill.background()

    def add_closing_slide(self, name: str = "", title: str = "",
                          contact: str = "", dark: bool = True):
        """Closing slide: Thank You (dark by default) + Oracle logo closing.

        Uses the FY26 Thank You layout (PH0 = "Thank you", PH37 = name/title,
        PH36 = contact/URL), followed by the Oracle Closing Slide layout
        (logo-only, no placeholders).

        Args:
            name:    Presenter name, e.g. "Diego Cabrera"
            title:   Presenter title/firm, e.g. "Solutions Architect | Oracle"
            contact: Contact line, e.g. email or URL
            dark:    Use dark variant (default True)
        """
        layout_idx = Layouts.THANK_YOU_DARK if dark else Layouts.THANK_YOU_LIGHT
        slide = self._add_layout_slide(layout_idx)
        # PH0 defaults to "Thank you" from the template — clear if not wanted
        # PH37: name + title on one line
        if name or title:
            parts = [p for p in [name, title] if p]
            self._set_placeholder(slide, 37, "  |  ".join(parts))
        # PH36: contact / URL
        if contact:
            self._set_placeholder(slide, 36, contact)

        # Oracle Closing Slide (logo-only)
        closing_idx = Layouts.CLOSING_DARK if dark else Layouts.CLOSING_LIGHT
        self._add_layout_slide(closing_idx)

    def save(self, filepath: str):
        """Save the presentation to a .pptx file."""
        self.prs.save(filepath)

    @property
    def slide_count(self):
        return len(self.prs.slides)
