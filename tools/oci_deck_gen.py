#!/usr/bin/env python3
"""
OCI Deal Accelerator — Slide Deck Generator (.pptx)

Produces a 10-12 slide architecture proposal deck using the Oracle FY26 official
PowerPoint template and the Oracle Redwood design system colors and typography.

Usage:
    python oci_deck_gen.py --spec proposal-data.yaml --output proposal.pptx

Or import and use programmatically:
    from oci_deck_gen import OCIDeckGenerator
    gen = OCIDeckGenerator(customer="Acme Corp", project="DB Migration")
    gen.add_summary_slide(...)
    gen.save("proposal.pptx")
"""

import yaml
import argparse
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from oci_pptx_base import Colors, Layouts, OraclePresBase


# ============================================================
# Defensive input helpers
# ============================================================

def pick(mapping: dict, *keys, default=""):
    """Return the first non-empty value for any of the provided keys."""
    if not isinstance(mapping, dict):
        return default
    for key in keys:
        if key in mapping:
            value = mapping.get(key)
            if value is not None:
                return value
    return default


def coerce_list(value):
    """Accept list, string, tuple, or None and normalize to a list."""
    if value is None:
        return []
    if isinstance(value, list):
        return value
    if isinstance(value, tuple):
        return list(value)
    if isinstance(value, str):
        return [value] if value else []
    return [value]


def pick_list(mapping: dict, *keys):
    """Return the first list-like field normalized to a list."""
    if not isinstance(mapping, dict):
        return []
    for key in keys:
        if key in mapping:
            return coerce_list(mapping.get(key))
    return []


# ============================================================
# MCP flat-spec adapter
# ============================================================

def _is_flat_spec(spec: dict) -> bool:
    """MCP payload shape: no 'metadata' key, but has customer_name/id/title."""
    if not isinstance(spec, dict):
        return False
    if "metadata" in spec:
        return False
    return any(k in spec for k in ("customer_name", "customer_id", "title"))


_PATTERNS_DIR = Path(__file__).resolve().parent.parent / "kb" / "patterns"


def _load_patterns_yaml(filename: str) -> dict:
    """Load a patterns YAML (supports 2-document layout: header + body)."""
    path = _PATTERNS_DIR / filename
    if not path.is_file():
        return {}
    try:
        with open(path, "r") as f:
            docs = [d for d in yaml.safe_load_all(f) if isinstance(d, dict)]
        return docs[-1] if docs else {}
    except Exception:
        return {}


def _default_architecture_principles() -> dict:
    """Curated ECAL principles — 3 per category, 'applies_when: always' first."""
    data = _load_patterns_yaml("architecture-principles.yaml")
    principles = data.get("principles", {}) if isinstance(data, dict) else {}
    result = {}
    for cat in ("design", "deployment", "service"):
        items = principles.get(cat, []) or []
        always = [p for p in items if p.get("applies_when") == "always"]
        picks = (always or items)[:3]
        result[cat] = [
            {
                "id": p.get("id", ""),
                "name": p.get("name", ""),
                "summary": p.get("principle", p.get("summary", "")),
            }
            for p in picks
            if isinstance(p, dict)
        ]
    return result if any(result.values()) else {}


def _default_operational_raci(model: str = "co_managed") -> list:
    data = _load_patterns_yaml("operational-raci.yaml")
    models = data.get("models", {}) if isinstance(data, dict) else {}
    chosen = models.get(model) or models.get("co_managed") or {}
    raci = chosen.get("raci", [])
    return [r for r in raci if isinstance(r, dict)]


def _derive_environments(deployment: str, dr: str) -> list:
    """Produce a default environment catalogue from deployment model + DR flag."""
    deployment_lc = (deployment or "").lower()
    dr_enabled = bool(dr) and str(dr).lower() not in {"", "no", "none", "not in scope", "disabled", "false"}
    envs = [
        {"name": "Production", "sizing": "Full capacity per workload profile",
         "isolation": "Dedicated compartment, private endpoints", "cost_pct": 55},
        {"name": "Pre-Production / UAT", "sizing": "50-75% of production",
         "isolation": "Shared tenancy OK, same region", "cost_pct": 20},
        {"name": "Dev/Test", "sizing": "25% of production, auto-stop off-hours",
         "isolation": "Shared, relaxed controls", "cost_pct": 10},
    ]
    if dr_enabled or "multi" in deployment_lc or "dr" in deployment_lc:
        envs.append({
            "name": "DR (secondary region)",
            "sizing": "100% of production, async replication",
            "isolation": "Cross-region, same security baseline",
            "cost_pct": 15,
        })
    return envs


_HA_DR_BY_TIER = {
    "platinum": {
        "technology": "RAC + cross-region Active Data Guard (sync)",
        "rto": "< 1 hour",
        "rpo": "< 5 minutes",
    },
    "gold": {
        "technology": "TAC + cross-region Active Data Guard (async)",
        "rto": "2-4 hours",
        "rpo": "15 minutes",
    },
    "silver": {
        "technology": "Single instance + cross-region backup restore",
        "rto": "4-8 hours",
        "rpo": "1 hour",
    },
    "bronze": {
        "technology": "Single instance + Object Storage archive",
        "rto": "24 hours",
        "rpo": "24 hours",
    },
}


def _derive_ha_dr_tiers(workloads: list) -> list:
    """Produce one HA/DR row per distinct service tier present."""
    seen = []
    for wl in workloads or []:
        tier = str(wl.get("tier", "")).lower()
        if tier in _HA_DR_BY_TIER and tier not in seen:
            seen.append(tier)
    return [
        {"tier": t.title(), **_HA_DR_BY_TIER[t]}
        for t in seen
    ]


def _default_security_baseline() -> dict:
    """Conservative OCI security baseline — customer confirms specifics."""
    return {
        "controls": {
            "identity": [
                "IAM compartments per environment + tag-based policies",
                "Federated SSO via customer IdP (SAML/OIDC)",
                "MFA enforced for all human users",
            ],
            "network": [
                "Private subnets with NSGs; no public IPs on data tier",
                "Service Gateway for OCI APIs; no egress over public internet",
                "FastConnect or Site-to-Site VPN for hybrid connectivity",
            ],
            "database": [
                "TDE at rest (customer-managed keys via OCI Vault)",
                "Network encryption in transit (TLS 1.2+)",
                "Data Safe for activity auditing + sensitive data discovery",
            ],
            "monitoring": [
                "OCI Logging + Service Connector Hub to SIEM",
                "Cloud Guard detectors enabled (baseline recipe)",
                "OCI Monitoring alarms on quota, cost, and security events",
            ],
        },
        "compliance": ["ISO 27001", "SOC 2", "PCI-DSS"],
    }


def _derive_service_tiering(services: list) -> list:
    """Map a flat services list to service-tiering workload rows.

    Preserves any tier/uptime/rto/rpo the caller supplied; fills sensible
    defaults so the slide always has a full table instead of blanks.
    """
    workloads = []
    for svc in services or []:
        if not isinstance(svc, dict):
            continue
        name = svc.get("name") or svc.get("workload") or svc.get("sku") or "Service"
        tier = svc.get("tier") or svc.get("service_tier") or "Gold"
        workloads.append({
            "name": name,
            "tier": str(tier).title(),
            "uptime": svc.get("uptime", "99.9%"),
            "rto": svc.get("rto", "4 hours"),
            "rpo": svc.get("rpo", "15 minutes"),
        })
    return workloads


def _adapt_flat_spec(spec: dict) -> dict:
    """Map the MCP flat payload to the proposal-spec structure from_spec expects."""
    arch = spec.get("architecture") if isinstance(spec.get("architecture"), dict) else {}
    tenancy = spec.get("tenancy") if isinstance(spec.get("tenancy"), dict) else {}

    customer = spec.get("customer_name") or spec.get("customer_id", "")
    title = spec.get("title") or arch.get("workload", "") or "Architecture Proposal"
    workload = spec.get("workload_type") or arch.get("workload", "")
    deployment = spec.get("deployment_model") or arch.get("deployment", "")
    billing = spec.get("billing_model") or arch.get("license_model", "")
    capacity = spec.get("capacity") or arch.get("capacity") or {}
    services = spec.get("services", []) or []

    # Region discovery — primary + DR can arrive as flat fields, under
    # architecture.*, or as a tenancy.regions list (items may be plain
    # strings or {region, role} dicts).
    primary_region = (
        spec.get("primary_region")
        or spec.get("region")
        or arch.get("primary_region")
        or arch.get("region")
        or ""
    )
    dr_region = (
        spec.get("dr_region")
        or spec.get("secondary_region")
        or arch.get("dr_region")
        or arch.get("secondary_region")
        or ""
    )
    tenancy_regions = tenancy.get("regions") or spec.get("regions") or arch.get("regions") or []
    if isinstance(tenancy_regions, list):
        for entry in tenancy_regions:
            if isinstance(entry, dict):
                name = entry.get("name") or entry.get("region") or entry.get("id") or ""
                role = str(entry.get("role") or entry.get("type") or "").lower()
                if role in ("primary", "hub", "prod", "production") and not primary_region:
                    primary_region = name
                elif role in ("dr", "standby", "secondary", "failover") and not dr_region:
                    dr_region = name
            elif isinstance(entry, str):
                if not primary_region:
                    primary_region = entry
                elif not dr_region:
                    dr_region = entry
    region = primary_region

    dr_raw = spec.get("disaster_recovery")
    if dr_raw is None and "dr" in arch:
        dr_raw = "Enabled" if arch.get("dr") else "Not in scope"
    # If the spec names a DR region, treat DR as enabled even when the
    # disaster_recovery flag is omitted.
    if not dr_raw and dr_region:
        dr_raw = f"Enabled — secondary region: {dr_region}"

    target_parts = [p for p in (workload, deployment and f"deployment: {deployment}",
                                region and f"region: {region}") if p]
    target_state = " — ".join(target_parts)

    current_state = []
    if isinstance(capacity, dict):
        for k, v in capacity.items():
            current_state.append(f"{k.replace('_', ' ').title()}: {v}")
    if billing:
        current_state.append(f"License model: {billing}")
    if dr_raw:
        current_state.append(f"Disaster recovery: {dr_raw}")

    line_items = []
    for svc in services:
        if not isinstance(svc, dict):
            continue
        qty = svc.get("quantity", svc.get("qty"))
        notes = svc.get("notes", "")
        if qty is not None:
            notes = f"Qty: {qty}" + (f" — {notes}" if notes else "")
        line_items.append({"component": svc.get("name", ""), "monthly_payg": "—", "notes": notes})

    cost_summary = spec.get("cost_summary") or {}
    monthly = cost_summary.get("monthly_estimate")
    if monthly is not None:
        currency = cost_summary.get("currency", "USD")
        annual = cost_summary.get("annual_estimate", monthly * 12)
        line_items.append({
            "component": "Total (estimated)",
            "monthly_payg": f"{currency} {monthly:,.2f}",
            "notes": f"Annual: {currency} {annual:,.2f}",
        })

    adapted = {
        "metadata": {"customer": customer, "subtitle": title},
        "summary": {
            "why": f"{customer} — {title}" if customer and title else (customer or title),
            "current_state": current_state,
            "target_state": target_state,
            "timeline": "",
        },
    }

    workloads = _derive_service_tiering(services)
    if workloads:
        adapted["service_tiering"] = workloads

    arch_desc_parts = [p for p in (
        workload and f"Workload: {workload}",
        deployment and f"Deployment: {deployment}",
        region and f"Primary region: {region}",
        dr_region and f"DR region: {dr_region}",
        billing and f"Licensing: {billing}",
    ) if p]
    architecture_block: dict = {}
    # Preserve any caller-supplied visual/diagram_path instead of overwriting.
    if isinstance(arch.get("visual"), dict):
        architecture_block["visual"] = arch["visual"]
    if arch.get("diagram_path"):
        architecture_block["diagram_path"] = arch["diagram_path"]
    if arch_desc_parts:
        architecture_block["description"] = "  •  ".join(arch_desc_parts)
    # When we know both regions and no visual was supplied, emit a simple
    # two-region visual so the DR region renders with its real name.
    if region and dr_region and "visual" not in architecture_block:
        architecture_block["visual"] = {
            "regions": [
                {"name": region, "primary": True, "label": "PRIMARY"},
                {"name": dr_region, "primary": False, "label": "DR STANDBY"},
            ],
        }
    if architecture_block:
        adapted["architecture"] = architecture_block

    principles = _default_architecture_principles()
    if principles:
        adapted["architecture_principles"] = principles

    ha_dr_tiers = _derive_ha_dr_tiers(workloads)
    if ha_dr_tiers:
        adapted["ha_dr"] = {
            "tiers": ha_dr_tiers,
            "description": "HA/DR posture per service tier. Confirm RTO/RPO with business stakeholders.",
        }

    adapted["security"] = _default_security_baseline()

    environments = _derive_environments(deployment, dr_raw)
    if environments:
        adapted["environment_catalogue"] = {"environments": environments}

    # Respect user-supplied RACI if present (string model name, list of items,
    # or dict with model/raci_items); otherwise fall back to co_managed default.
    user_raci = None
    for key in ("operational_raci", "raci", "operations_raci"):
        if spec.get(key) is not None:
            user_raci = spec[key]
            break
    if user_raci is not None:
        adapted["operational_raci"] = user_raci
    else:
        raci_items = _default_operational_raci("co_managed")
        if raci_items:
            adapted["operational_raci"] = {
                "raci_items": raci_items,
                "model": "co_managed",
            }

    adapted["next_steps"] = {
        "steps": [
            "Validate architecture + service tiers with technical stakeholders",
            "Confirm capacity assumptions against current/expected utilization",
            "Lock commitment model (PAYG vs UCM) and finalize discount",
            "Define proof-of-concept / pilot scope and success criteria",
            "Schedule operational handover and RACI walkthrough",
        ],
    }

    if line_items:
        adapted["cost"] = {
            "line_items": line_items,
            "assumptions": spec.get("assumptions", []) or [],
            "show_byol": False,
        }
    return adapted


# ============================================================
# Slide Generator
# ============================================================

class OCIDeckGenerator(OraclePresBase):
    """Generate Oracle Redwood-styled architecture proposal slide decks."""

    SLIDE_WIDTH = Inches(13.333)   # Widescreen 16:9
    SLIDE_HEIGHT = Inches(7.5)
    MARGIN = Inches(0.5)

    # Use Oracle Red accent for proposal deck (differentiates from bizcase teal)
    TITLE_ACCENT_COLOR = Colors.ORACLE_RED

    def __init__(self, customer: str = "", project: str = "",
                 architect: str = "", firm: str = "",
                 template: Optional[str] = None):
        super().__init__(template)
        self.customer = customer
        self.project = project
        self.architect = architect
        self.firm = firm

    # ---- Architecture visual helpers ----

    def _add_service_block(self, slide, left, top, width, height,
                           label, fill_color):
        """Add a colored service block with white text label."""
        shape = self._add_rect(
            slide, left, top, width, height, fill_color=fill_color,
        )
        shape.text_frame.word_wrap = True
        p = shape.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = Colors.WHITE
        p.font.name = self.FONT
        p.alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].space_before = Pt(0)
        shape.text_frame.paragraphs[0].space_after = Pt(0)
        return shape

    def _add_container(self, slide, left, top, width, height,
                       label, border_color, fill_color=None,
                       dashed=True, label_size=10):
        """Add a container outline (region, VCN, subnet)."""
        shape = self._add_rect(
            slide, left, top, width, height,
            fill_color=fill_color, border_color=border_color,
        )
        if not fill_color:
            shape.fill.background()
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        if dashed:
            shape.line.dash_style = 2  # MSO_LINE_DASH_STYLE.DASH
        # Label at top-left inside container
        self._add_textbox(
            slide, left + Inches(0.1), top + Inches(0.05),
            width - Inches(0.2), Inches(0.3),
            text=label, font_size=label_size, bold=True,
            color=border_color,
        )
        return shape

    def _add_arrow(self, slide, start_x, start_y, end_x, end_y,
                   color=None, dashed=False, label=""):
        """Add a connector arrow between two points."""
        connector = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            start_x, start_y, end_x, end_y,
        )
        connector.line.color.rgb = color or RGBColor(0x70, 0x6E, 0x6F)
        connector.line.width = Pt(1.5)
        if dashed:
            connector.line.dash_style = 2
        return connector

    # ---- Slide Methods ----

    def add_title_slide(self, subtitle: str = ""):
        """Slide 1: Title slide using FY26 Dark Cover layout."""
        slide = self._add_layout_slide(Layouts.COVER_DARK)

        customer = self.customer or "Architecture Proposal"
        date_str = datetime.now().strftime("%B %Y")
        sub = subtitle or f"Architecture Proposal — {date_str}"

        self._set_placeholder(slide, 0, customer)
        self._set_placeholder(slide, 33, sub)
        self._set_placeholder(slide, 35, date_str)
        if self.architect or self.firm:
            info = self.architect
            if self.firm:
                info += f"  |  {self.firm}" if self.architect else self.firm
            self._set_placeholder(slide, 34, info)

    def add_summary_slide(self, why: str, current_state: list,
                          target_state: str, timeline: str):
        """Slide 2: Engagement Summary."""
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Engagement Summary", margin=self.MARGIN)

        y = Inches(1.2)

        # Why — teal for emphasis
        self._add_textbox(
            slide, self.MARGIN, y, Inches(11), Inches(0.5),
            text=why, font_size=14, color=Colors.TEAL, bold=True,
            font_name=self.FONT_HEADING,
        )
        y += Inches(0.7)

        # Current state
        self._add_textbox(
            slide, self.MARGIN, y, Inches(5), Inches(0.4),
            text="Current State", font_size=14, bold=True,
        )
        y += Inches(0.45)

        for item in current_state:
            self._add_textbox(
                slide, Inches(0.7), y, Inches(10), Inches(0.35),
                text=f"•  {item}", font_size=12,
            )
            y += Inches(0.35)

        y += Inches(0.3)

        # Target state
        self._add_textbox(
            slide, self.MARGIN, y, Inches(5), Inches(0.4),
            text="Target State", font_size=14, bold=True,
        )
        y += Inches(0.45)
        self._add_textbox(
            slide, Inches(0.7), y, Inches(10), Inches(0.4),
            text=target_state, font_size=12, color=Colors.TEAL,
        )
        y += Inches(0.5)

        # Timeline
        self._add_textbox(
            slide, self.MARGIN, y, Inches(10), Inches(0.4),
            text=f"Timeline: {timeline}", font_size=12, bold=True,
        )

    def add_architecture_slide(self, diagram_path: Optional[str] = None,
                                description: str = "",
                                visual: Optional[dict] = None):
        """Slide 3: Architecture Overview with diagram or visual layout.

        visual: optional dict with structured architecture data for rendering
                as colored blocks. Keys: regions, on_prem, security_footer.
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Architecture Overview", margin=self.MARGIN)

        if diagram_path:
            try:
                slide.shapes.add_picture(
                    diagram_path,
                    Inches(0.5), Inches(1.1),
                    Inches(12.3), Inches(6),
                )
            except Exception:
                self._add_textbox(
                    slide, Inches(1), Inches(2.5),
                    Inches(10), Inches(2),
                    text=f"[Diagram: {diagram_path}]\n\n{description}",
                    font_size=14, color=Colors.SECONDARY_TEXT,
                )
        elif visual:
            self._render_architecture_visual(slide, visual)
        else:
            self._add_textbox(
                slide, Inches(1), Inches(2),
                Inches(11), Inches(4),
                text=description or "[Insert architecture diagram — export from .drawio file]",
                font_size=14, color=Colors.SECONDARY_TEXT, italic=True,
                alignment=PP_ALIGN.CENTER,
            )

    def _render_architecture_visual(self, slide, visual: dict):
        """Render a visual architecture diagram from structured data."""
        regions = visual.get("regions", [])
        on_prem = visual.get("on_prem")
        security = visual.get("security_footer", "")

        # Category colors for service blocks
        cat_colors = {
            "infrastructure": Colors.TEAL,
            "database": Colors.COPPER,
            "integration": Colors.PURPLE,
            "security": Colors.FOREST,
            "dormant": RGBColor(0xDF, 0xDC, 0xD8),
        }

        y_cursor = Inches(1.15)

        for reg_idx, region in enumerate(regions):
            is_primary = region.get("primary", reg_idx == 0)
            region_name = region.get("name", f"Region {reg_idx+1}")
            region_label = region.get("label", "PRIMARY" if is_primary else "DR STANDBY")

            # Region dimensions — primary gets more space
            if is_primary:
                reg_left = Inches(0.4)
                reg_width = Inches(8.2)
                reg_height = Inches(4.8)
                reg_top = y_cursor
            else:
                reg_left = Inches(8.9)
                reg_width = Inches(4.1)
                reg_height = Inches(4.8)
                reg_top = y_cursor

            # Region container — solid border, light fill
            self._add_container(
                slide, reg_left, reg_top, reg_width, reg_height,
                label=f"{region_name}  [{region_label}]",
                border_color=RGBColor(0x9E, 0x98, 0x92),
                fill_color=RGBColor(0xF5, 0xF4, 0xF2),
                dashed=False, label_size=10,
            )

            inner_left = reg_left + Inches(0.15)
            inner_top = reg_top + Inches(0.45)
            inner_width = reg_width - Inches(0.3)

            # VCN container if present
            vcn = region.get("vcn")
            if vcn:
                vcn_name = vcn.get("name", "VCN")
                vcn_cidr = vcn.get("cidr", "")
                vcn_label = f"{vcn_name} {vcn_cidr}" if vcn_cidr else vcn_name
                vcn_height = reg_height - Inches(0.6)

                self._add_container(
                    slide, inner_left, inner_top, inner_width, vcn_height,
                    label=vcn_label,
                    border_color=Colors.BURNT_ORANGE,
                    dashed=True, label_size=9,
                )

                # Render subnets
                subnets = vcn.get("subnets", [])
                sub_top = inner_top + Inches(0.4)
                sub_left = inner_left + Inches(0.15)
                sub_width = inner_width - Inches(0.3)

                for sub_idx, subnet in enumerate(subnets):
                    sub_name = subnet.get("name", f"Subnet {sub_idx+1}")
                    sub_border = Colors.BURNT_ORANGE
                    sub_h = Inches(0.95) if is_primary else Inches(0.7)

                    self._add_container(
                        slide, sub_left, sub_top, sub_width, sub_h,
                        label=sub_name,
                        border_color=sub_border,
                        fill_color=Colors.WARM_BG,
                        dashed=True, label_size=8,
                    )

                    # Service blocks inside subnet
                    services = subnet.get("services", [])
                    svc_left = sub_left + Inches(0.1)
                    svc_top = sub_top + Inches(0.3)
                    svc_spacing = Inches(0.08)
                    # Calculate service block width to fit
                    if services:
                        avail_width = sub_width - Inches(0.2)
                        svc_w = min(
                            Inches(2.2),
                            (avail_width - svc_spacing * (len(services) - 1)) / len(services)
                        )
                        svc_h = Inches(0.5)
                        for svc in services:
                            svc_name = svc.get("name", "Service")
                            svc_cat = svc.get("category", "infrastructure")
                            svc_color = cat_colors.get(svc_cat, Colors.TEAL)
                            self._add_service_block(
                                slide, svc_left, svc_top, svc_w, svc_h,
                                label=svc_name, fill_color=svc_color,
                            )
                            svc_left += svc_w + svc_spacing

                    sub_top += sub_h + Inches(0.1)

                # Gateways bar at bottom of VCN
                gateways = vcn.get("gateways", [])
                if gateways:
                    gw_top = inner_top + vcn_height - Inches(0.55)
                    gw_left = inner_left + Inches(0.15)
                    gw_text = "  |  ".join(gateways)
                    self._add_textbox(
                        slide, gw_left, gw_top,
                        inner_width - Inches(0.3), Inches(0.35),
                        text=gw_text, font_size=8, bold=True,
                        color=Colors.TEAL,
                        alignment=PP_ALIGN.CENTER,
                    )

            else:
                # No VCN — render services directly (e.g. DR standby)
                services = region.get("services", [])
                svc_top = inner_top + Inches(0.3)
                for svc in services:
                    svc_name = svc.get("name", "Service")
                    svc_cat = svc.get("category", "database")
                    svc_color = cat_colors.get(svc_cat, Colors.TEAL)
                    self._add_service_block(
                        slide, inner_left + Inches(0.3), svc_top,
                        inner_width - Inches(0.6), Inches(0.6),
                        label=svc_name, fill_color=svc_color,
                    )
                    svc_top += Inches(0.75)

                # DR details text
                dr_details = region.get("details", "")
                if dr_details:
                    self._add_textbox(
                        slide, inner_left + Inches(0.1), svc_top + Inches(0.1),
                        inner_width - Inches(0.2), Inches(0.6),
                        text=dr_details, font_size=8, italic=True,
                        color=Colors.SECONDARY_TEXT,
                        alignment=PP_ALIGN.CENTER,
                    )

        # On-prem block
        if on_prem:
            op_left = Inches(0.4)
            op_top = Inches(6.15)
            op_width = Inches(3.5)
            self._add_rect(
                slide, op_left, op_top, op_width, Inches(0.45),
                fill_color=RGBColor(0x70, 0x66, 0x5E),
            )
            op_shape = slide.shapes[-1]
            p = op_shape.text_frame.paragraphs[0]
            p.text = on_prem.get("name", "On-Premises")
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = Colors.WHITE
            p.font.name = self.FONT
            p.alignment = PP_ALIGN.CENTER

            # Connection label
            conn_label = on_prem.get("connection", "IPSec VPN")
            self._add_textbox(
                slide, op_left + op_width + Inches(0.15), op_top,
                Inches(3), Inches(0.45),
                text=f"--- {conn_label} ---",
                font_size=8, italic=True, color=Colors.SECONDARY_TEXT,
                alignment=PP_ALIGN.LEFT,
            )

        # Security footer
        if security:
            self._add_textbox(
                slide, Inches(0.4), Inches(6.7),
                Inches(12.5), Inches(0.35),
                text=security, font_size=8, bold=True,
                color=Colors.TEAL,
                alignment=PP_ALIGN.LEFT,
            )

    def add_service_tiering_slide(self, workloads: list):
        """Service Tiering slide — maps workloads to Platinum/Gold/Silver/Bronze.

        workloads: list of {"name": str, "tier": str, "uptime": str,
                           "rto": str, "rpo": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Service Tiering", margin=self.MARGIN)

        # Subtitle
        self._add_textbox(
            slide, self.MARGIN, Inches(1.1),
            Inches(12), Inches(0.4),
            text="Each tier drives: HA/DR topology, backup strategy, isolation model, support level.",
            font_size=11, italic=True, color=Colors.SECONDARY_TEXT,
        )

        rows = len(workloads) + 1
        table = self._add_table(
            slide, rows, 5,
            self.MARGIN, Inches(1.7),
            Inches(12), Inches(0.45 * rows),
        )

        table.columns[0].width = Inches(3.0)
        table.columns[1].width = Inches(2.5)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(2.25)
        table.columns[4].width = Inches(2.25)

        headers = ["Workload", "Tier", "Uptime", "RTO", "RPO"]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER,
            )

        tier_colors = {
            "platinum": Colors.ORACLE_RED,
            "gold": Colors.BURNT_ORANGE,
            "silver": Colors.SECONDARY_TEXT,
            "bronze": Colors.MUTED_TEAL,
        }

        for i, wl in enumerate(workloads):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            workload_name = pick(wl, "name", "workload", "title")
            self._style_table_cell(table.cell(row_idx, 0), workload_name, font_size=10, bold=True, bg_color=bg)
            tier_label = pick(wl, "tier", "name", default="Silver")
            tier_color = tier_colors.get(tier_label.lower(), Colors.SECONDARY_TEXT)
            self._style_table_cell(
                table.cell(row_idx, 1), tier_label.title(),
                font_size=10, bold=True, color=tier_color, bg_color=bg,
                alignment=PP_ALIGN.CENTER,
            )
            self._style_table_cell(table.cell(row_idx, 2), wl.get("uptime", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
            self._style_table_cell(table.cell(row_idx, 3), wl.get("rto", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
            self._style_table_cell(table.cell(row_idx, 4), wl.get("rpo", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)

    def add_architecture_principles_slide(self, principles: dict):
        """Architecture Principles slide — ECAL Design/Deployment/Service categories.

        principles: {"design": [{"id": str, "name": str, "summary": str}],
                     "deployment": [...], "service": [...]}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Architecture Principles", margin=self.MARGIN)

        y = Inches(1.2)
        col_x = self.MARGIN
        col_width = Inches(4)

        for category in ["design", "deployment", "service"]:
            items = principles.get(category, [])
            if not items:
                continue

            # Category heading
            self._add_textbox(
                slide, col_x, y, col_width, Inches(0.35),
                text=category.upper(), font_size=13, bold=True,
                color=Colors.TEAL,
            )

            item_y = y + Inches(0.4)
            for item in items:
                pid = item.get("id", "")
                name = item.get("name", "")
                summary = item.get("summary", "")
                label = f"{pid}  {name}" if pid else name
                if summary:
                    label += f" — {summary}"
                self._add_textbox(
                    slide, col_x + Inches(0.1), item_y,
                    col_width - Inches(0.1), Inches(0.4),
                    text=f"• {label}", font_size=11,
                )
                item_y += Inches(0.4)

            col_x += Inches(4.2)

    def add_decisions_slide(self, decisions: list):
        """Slide 4: Architecture Decisions table.

        decisions: list of {"decision": str, "rationale": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Architecture Decisions", margin=self.MARGIN)

        rows = len(decisions) + 1
        table = self._add_table(
            slide, rows, 2,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.5 * rows),
        )

        # Set column widths
        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(7.5)

        # Header row
        self._style_table_cell(
            table.cell(0, 0), "Decision", font_size=11, bold=True,
            color=Colors.WHITE, bg_color=Colors.TEAL,
        )
        self._style_table_cell(
            table.cell(0, 1), "Rationale", font_size=11, bold=True,
            color=Colors.WHITE, bg_color=Colors.TEAL,
        )

        # Data rows
        for i, dec in enumerate(decisions):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(
                table.cell(row_idx, 0), pick(dec, "decision", "title"),
                font_size=11, bold=True, bg_color=bg,
            )
            self._style_table_cell(
                table.cell(row_idx, 1), pick(dec, "rationale", "reason", "notes"),
                font_size=10, bg_color=bg,
            )

    def add_ha_dr_slide(self, tiers: list, description: str = ""):
        """Slide 5: HA/DR table.

        tiers: list of {"tier": str, "technology": str, "rto": str, "rpo": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "High Availability & Disaster Recovery", margin=self.MARGIN)

        if description:
            self._add_textbox(
                slide, self.MARGIN, Inches(1.2),
                Inches(12), Inches(0.5),
                text=description, font_size=12,
            )

        rows = len(tiers) + 1
        table = self._add_table(
            slide, rows, 4,
            self.MARGIN, Inches(1.9),
            Inches(12), Inches(0.45 * rows),
        )

        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(4.5)
        table.columns[2].width = Inches(2.5)
        table.columns[3].width = Inches(2.5)

        headers = ["Tier", "Technology", "RTO", "RPO"]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER,
            )

        for i, tier in enumerate(tiers):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(table.cell(row_idx, 0), pick(tier, "tier", "name"), font_size=10, bold=True, bg_color=bg)
            self._style_table_cell(table.cell(row_idx, 1), pick(tier, "technology", "architecture", "pattern"), font_size=10, bg_color=bg)
            self._style_table_cell(table.cell(row_idx, 2), pick(tier, "rto"), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
            self._style_table_cell(table.cell(row_idx, 3), pick(tier, "rpo"), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)

    def add_security_slide(self, controls: dict, compliance: list = None):
        """Slide 6: Security & Compliance.

        controls: {"identity": [...], "network": [...], "database": [...], "monitoring": [...]}
        compliance: ["PCI-DSS", "SOC2", ...] or None
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Security & Compliance", margin=self.MARGIN)

        y = Inches(1.2)

        # Compliance badges
        if compliance:
            x = self.MARGIN
            for framework in compliance:
                shape = self._add_rect(
                    slide, x, y, Inches(1.5), Inches(0.4),
                    fill_color=Colors.SUCCESS,
                )
                shape.text_frame.paragraphs[0].text = f"✓ {framework}"
                shape.text_frame.paragraphs[0].font.size = Pt(10)
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
                shape.text_frame.paragraphs[0].font.name = self.FONT
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                x += Inches(1.7)
            y += Inches(0.7)

        # Controls in columns
        col_x = self.MARGIN
        col_width = Inches(3)
        for area, items in controls.items():
            area_title = area.replace("_", " ").title()
            self._add_textbox(
                slide, col_x, y, col_width, Inches(0.35),
                text=area_title, font_size=12, bold=True,
                color=Colors.TEAL,
            )
            item_y = y + Inches(0.4)
            for item in items:
                self._add_textbox(
                    slide, col_x + Inches(0.1), item_y,
                    col_width - Inches(0.1), Inches(0.3),
                    text=f"• {item}", font_size=10,
                )
                item_y += Inches(0.3)
            col_x += Inches(3.2)

    def add_environment_catalogue_slide(self, environments: list,
                                       cost_notes=None):
        """Environment Catalogue slide — Prod/Pre-Prod/Dev-Test/DR per workload.

        environments: list of {"name"|"environment": str, "sizing": str,
                              "isolation": str, "cost_pct": int}
        cost_notes: optional list of strings or single string
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Environment Catalogue", margin=self.MARGIN)

        # Normalize cost_notes to a list
        if isinstance(cost_notes, str):
            cost_notes = [s.strip() for s in cost_notes.split(".") if s.strip()]

        rows = len(environments) + 1
        table = self._add_table(
            slide, rows, 5,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.4 * rows),
        )

        # Detect schema: new (name/sizing/isolation/cost_pct) vs legacy (environment/tier/databases/ocpus/isolation)
        sample = environments[0] if environments else {}
        use_new_schema = "sizing" in sample or "cost_pct" in sample

        if use_new_schema:
            table = self._add_table(
                slide, rows, 4,
                self.MARGIN, Inches(1.2),
                Inches(12), Inches(0.4 * rows),
            )
            table.columns[0].width = Inches(2.5)
            table.columns[1].width = Inches(5.0)
            table.columns[2].width = Inches(2.5)
            table.columns[3].width = Inches(2.0)

            headers = ["Environment", "Sizing", "Isolation", "Cost %"]
            for j, h in enumerate(headers):
                self._style_table_cell(
                    table.cell(0, j), h, font_size=11, bold=True,
                    color=Colors.WHITE, bg_color=Colors.TEAL,
                    alignment=PP_ALIGN.CENTER,
                )

            for i, env in enumerate(environments):
                row_idx = i + 1
                bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
                self._style_table_cell(table.cell(row_idx, 0), env.get("name", env.get("environment", "")), font_size=10, bold=True, bg_color=bg)
                self._style_table_cell(table.cell(row_idx, 1), env.get("sizing", ""), font_size=10, bg_color=bg)
                self._style_table_cell(table.cell(row_idx, 2), env.get("isolation", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
                cost_pct = env.get("cost_pct", "")
                self._style_table_cell(table.cell(row_idx, 3), f"{cost_pct}%" if cost_pct != "" else "", font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
        else:
            table = self._add_table(
                slide, rows, 5,
                self.MARGIN, Inches(1.2),
                Inches(12), Inches(0.4 * rows),
            )
            table.columns[0].width = Inches(2.2)
            table.columns[1].width = Inches(1.8)
            table.columns[2].width = Inches(2.5)
            table.columns[3].width = Inches(1.5)
            table.columns[4].width = Inches(4.0)

            headers = ["Environment", "Tier", "Databases", "OCPUs", "Isolation"]
            for j, h in enumerate(headers):
                self._style_table_cell(
                    table.cell(0, j), h, font_size=11, bold=True,
                    color=Colors.WHITE, bg_color=Colors.TEAL,
                    alignment=PP_ALIGN.CENTER,
                )

            for i, env in enumerate(environments):
                row_idx = i + 1
                bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
                self._style_table_cell(table.cell(row_idx, 0), env.get("environment", ""), font_size=10, bold=True, bg_color=bg)
                self._style_table_cell(table.cell(row_idx, 1), env.get("tier", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
                self._style_table_cell(table.cell(row_idx, 2), env.get("databases", ""), font_size=10, bg_color=bg)
                self._style_table_cell(table.cell(row_idx, 3), env.get("ocpus", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
                self._style_table_cell(table.cell(row_idx, 4), env.get("isolation", ""), font_size=10, bg_color=bg)

        if cost_notes:
            notes_y = Inches(1.2) + Inches(0.4 * rows) + Inches(0.3)
            self._add_textbox(
                slide, self.MARGIN, notes_y,
                Inches(12), Inches(0.35),
                text="Cost Optimization", font_size=12, bold=True,
                color=Colors.TEAL,
            )
            note_y = notes_y + Inches(0.35)
            for note in cost_notes:
                self._add_textbox(
                    slide, self.MARGIN + Inches(0.1), note_y,
                    Inches(11.5), Inches(0.35),
                    text=f"• {note}", font_size=11,
                )
                note_y += Inches(0.35)

    def add_cost_slide(self, line_items: list, assumptions: list = None,
                       show_byol: bool = True):
        """Slide 7: Cost Estimate.

        line_items: list of {"component": str, "monthly_payg": str,
                            "monthly_byol": str (optional), "notes": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Cost Estimate", margin=self.MARGIN)

        cols = 4 if show_byol else 3
        rows = len(line_items) + 1
        table = self._add_table(
            slide, rows, cols,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.4 * rows),
        )

        if show_byol:
            table.columns[0].width = Inches(3.5)
            table.columns[1].width = Inches(2.5)
            table.columns[2].width = Inches(2.5)
            table.columns[3].width = Inches(3.5)
            headers = ["Component", "Monthly (PAYG)", "Monthly (BYOL)", "Notes"]
        else:
            table.columns[0].width = Inches(4)
            table.columns[1].width = Inches(3)
            table.columns[2].width = Inches(5)
            headers = ["Component", "Monthly", "Notes"]

        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
            )

        for i, item in enumerate(line_items):
            row_idx = i + 1
            component = pick(item, "component", "name", "label")
            notes = pick(item, "notes", "note")
            monthly_payg = pick(item, "monthly_payg", "monthly")
            is_total = "total" in str(component).lower()
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 and not is_total else None
            if is_total:
                bg = Colors.TEAL

            self._style_table_cell(
                table.cell(row_idx, 0), component,
                font_size=10, bold=is_total,
                color=Colors.WHITE if is_total else None,
                bg_color=bg,
            )
            self._style_table_cell(
                table.cell(row_idx, 1), monthly_payg,
                font_size=10, bold=is_total,
                color=Colors.WHITE if is_total else None,
                bg_color=bg, alignment=PP_ALIGN.RIGHT,
            )
            if show_byol:
                self._style_table_cell(
                    table.cell(row_idx, 2), item.get("monthly_byol", ""),
                    font_size=10, bold=is_total,
                    color=Colors.WHITE if is_total else None,
                    bg_color=bg, alignment=PP_ALIGN.RIGHT,
                )
                self._style_table_cell(
                    table.cell(row_idx, 3), notes,
                    font_size=9, bg_color=bg,
                    color=Colors.WHITE if is_total else Colors.SECONDARY_TEXT,
                )
            else:
                self._style_table_cell(
                    table.cell(row_idx, 2), notes,
                    font_size=9, bg_color=bg,
                    color=Colors.WHITE if is_total else Colors.SECONDARY_TEXT,
                )

        # Assumptions
        if assumptions:
            table_bottom = Inches(1.2) + Inches(0.4 * rows) + Inches(0.3)
            self._add_textbox(
                slide, self.MARGIN, table_bottom,
                Inches(12), Inches(0.35),
                text="Assumptions:", font_size=11, bold=True,
                color=Colors.SECONDARY_TEXT, italic=True,
            )
            for idx, assumption in enumerate(assumptions):
                self._add_textbox(
                    slide, Inches(0.7), table_bottom + Inches(0.35 + idx * 0.3),
                    Inches(11), Inches(0.3),
                    text=f"• {assumption}", font_size=10, italic=True,
                    color=Colors.SECONDARY_TEXT,
                )

    def add_cost_comparison_slide(self, rows: list, title: str = "Cost Comparison",
                                  col_headers: list = None):
        """Slide 8: Cost Comparison (optional).

        rows: list of {"item": str, "current": str, "oci": str, "savings": str}
        col_headers: custom column headers (default: ["Component", "Current", "OCI", "Savings"])
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, title, margin=self.MARGIN)

        headers = col_headers or ["Component", "Current", "OCI", "Savings"]
        num_rows = len(rows) + 1
        table = self._add_table(
            slide, num_rows, len(headers),
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.45 * num_rows),
        )

        col_widths = [Inches(4), Inches(2.5), Inches(2.5), Inches(3)]
        for j, w in enumerate(col_widths[:len(headers)]):
            table.columns[j].width = w

        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
            )

        for i, row in enumerate(rows):
            row_idx = i + 1
            is_total = "total" in row.get("item", "").lower()
            bg = Colors.TEAL if is_total else (Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None)
            txt_color = Colors.WHITE if is_total else None

            self._style_table_cell(
                table.cell(row_idx, 0), row.get("item", ""),
                font_size=10, bold=is_total, color=txt_color, bg_color=bg,
            )
            self._style_table_cell(
                table.cell(row_idx, 1), row.get("current", ""),
                font_size=10, bold=is_total, color=txt_color, bg_color=bg,
                alignment=PP_ALIGN.RIGHT,
            )
            self._style_table_cell(
                table.cell(row_idx, 2), row.get("oci", ""),
                font_size=10, bold=is_total, color=txt_color, bg_color=bg,
                alignment=PP_ALIGN.RIGHT,
            )
            if len(headers) > 3:
                savings_color = Colors.SUCCESS if not is_total else txt_color
                self._style_table_cell(
                    table.cell(row_idx, 3), row.get("savings", ""),
                    font_size=10, bold=True, color=savings_color, bg_color=bg,
                    alignment=PP_ALIGN.CENTER,
                )

    def add_migration_slide(self, phases: list, tools: list = None,
                            downtime: str = ""):
        """Migration Approach slide.

        phases: list of {"name": str, "duration"|"weeks": str,
                        "milestones"|"tasks": [...]}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Migration Approach", margin=self.MARGIN)

        phase_colors = [Colors.TEAL, Colors.BURNT_ORANGE, Colors.FOREST, Colors.PURPLE]

        n = len(phases)
        # Scale spacing to fit all phases: leave room for tools/downtime at bottom
        avail = 5.0  # inches available for phases (from y=1.3 to ~6.3)
        step = min(1.1, avail / max(n, 1))

        y = Inches(1.3)
        bar_left = Inches(2.8)
        bar_width = Inches(9.2)

        for i, phase in enumerate(phases):
            color = phase_colors[i % len(phase_colors)]
            duration = pick(phase, "duration", "weeks")
            milestones = pick_list(phase, "milestones", "tasks")

            # Phase label (left of bar)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(2.3), Inches(0.4),
                text=pick(phase, "name", "title"), font_size=11, bold=True,
            )

            # Colored phase bar with duration
            bar = self._add_rect(
                slide, bar_left, y + Inches(0.02),
                bar_width * 0.9, Inches(0.36),
                fill_color=color,
            )
            bar.text_frame.paragraphs[0].text = duration
            bar.text_frame.paragraphs[0].font.size = Pt(11)
            bar.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            bar.text_frame.paragraphs[0].font.name = self.FONT
            bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Milestones/tasks below bar (compact)
            if milestones:
                tasks_text = "  •  ".join(str(m) for m in milestones[:4])
                self._add_textbox(
                    slide, bar_left, y + Inches(0.40),
                    bar_width, Inches(0.30),
                    text=tasks_text, font_size=9, italic=True,
                    color=Colors.SECONDARY_TEXT,
                )

            y += Inches(step)

        # Tools
        if tools:
            y = max(y, Inches(1.3 + avail)) + Inches(0.15)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(12), Inches(0.35),
                text=f"Migration Tools: {', '.join(tools[:4])}",
                font_size=10, bold=True, color=Colors.TEAL,
            )
            y += Inches(0.35)

        # Downtime approach
        if downtime:
            if not tools:
                y = max(y, Inches(1.3 + avail)) + Inches(0.15)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(12), Inches(0.35),
                text=f"Downtime Approach: {downtime}",
                font_size=10, color=Colors.PRIMARY_TEXT,
            )

    def add_operational_raci_slide(self, raci_items: list,
                                   model: str = "co_managed"):
        """Operational RACI slide — responsibility matrix.

        raci_items: list of {"activity": str, "customer": str, "oracle": str}
        model: "fully_managed", "co_managed", or "self_managed"
        """
        slide = self._add_blank_slide()
        model_label = model.replace("_", "-").title()
        self._add_title_bar(slide, f"Operational Responsibilities ({model_label})", margin=self.MARGIN)

        rows = len(raci_items) + 1
        table = self._add_table(
            slide, rows, 3,
            self.MARGIN, Inches(1.2),
            Inches(10), Inches(0.38 * rows),
        )

        table.columns[0].width = Inches(5.0)
        table.columns[1].width = Inches(2.5)
        table.columns[2].width = Inches(2.5)

        headers = ["Activity", "Customer", "Oracle / Partner"]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER,
            )

        for i, item in enumerate(raci_items):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(table.cell(row_idx, 0), item.get("activity", ""), font_size=10, bg_color=bg)
            self._style_table_cell(table.cell(row_idx, 1), item.get("customer", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
            self._style_table_cell(table.cell(row_idx, 2), item.get("oracle", ""), font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)

        # Legend
        legend_y = Inches(1.2) + Inches(0.38 * rows) + Inches(0.2)
        self._add_textbox(
            slide, self.MARGIN, legend_y,
            Inches(10), Inches(0.35),
            text="R = Responsible    A = Accountable    C = Consulted    I = Informed",
            font_size=11, italic=True, color=Colors.SECONDARY_TEXT,
        )

    def add_risk_slide(self, risks: list):
        """Slide 10: Risk Register.

        risks: list of {"risk": str, "severity": "HIGH"|"MEDIUM"|"LOW", "mitigation": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Risk Register", margin=self.MARGIN)

        rows = len(risks) + 1
        table = self._add_table(
            slide, rows, 3,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.5 * rows),
        )

        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(6)

        headers = ["Risk", "Severity", "Mitigation"]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT,
            )

        severity_colors = {
            "HIGH": Colors.ERROR,
            "MEDIUM": Colors.WARNING,
            "LOW": Colors.SUCCESS,
        }

        for i, risk in enumerate(risks):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(
                table.cell(row_idx, 0), pick(risk, "risk", "title"),
                font_size=10, bg_color=bg,
            )
            sev = risk.get("severity", "MEDIUM").upper()
            self._style_table_cell(
                table.cell(row_idx, 1), sev,
                font_size=10, bold=True,
                color=severity_colors.get(sev, Colors.WARNING),
                bg_color=bg, alignment=PP_ALIGN.CENTER,
            )
            self._style_table_cell(
                table.cell(row_idx, 2), pick(risk, "mitigation", "action", "response"),
                font_size=10, bg_color=bg,
            )

    def add_scorecard_slide(self, pillars: list, recommendations: list = None):
        """Slide 11: Well-Architected Scorecard.

        pillars: list of {"name": str, "status": "PASS"|"PASS_WITH_RECOMMENDATIONS"|"GAPS_IDENTIFIED"|"NOT_APPLICABLE",
                         "passed": int, "total": int}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Well-Architected Scorecard", margin=self.MARGIN)

        status_styles = {
            "PASS": {"icon": "✓", "color": Colors.SUCCESS, "label": "PASS"},
            "PASS_WITH_RECOMMENDATIONS": {"icon": "⚠", "color": Colors.WARNING, "label": "PASS WITH RECOMMENDATIONS"},
            "GAPS_IDENTIFIED": {"icon": "✗", "color": Colors.ERROR, "label": "GAPS IDENTIFIED"},
            "NOT_APPLICABLE": {"icon": "—", "color": Colors.SECONDARY_TEXT, "label": "N/A"},
        }

        y = Inches(1.4)
        for pillar in pillars:
            status = pick(pillar, "status", default="PASS")
            style = status_styles.get(status, status_styles["PASS"])

            # Status indicator pill
            pill = self._add_rect(
                slide, self.MARGIN, y,
                Inches(0.5), Inches(0.45),
                fill_color=style["color"],
            )
            pill.text_frame.paragraphs[0].text = style["icon"]
            pill.text_frame.paragraphs[0].font.size = Pt(14)
            pill.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            pill.text_frame.paragraphs[0].font.name = self.FONT
            pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Pillar name
            self._add_textbox(
                slide, Inches(1.2), y,
                Inches(4), Inches(0.45),
                text=pick(pillar, "name", "pillar"), font_size=13, bold=True,
            )

            # Score
            passed = pick(pillar, "passed", default=0)
            total = pick(pillar, "total", default=0)
            score_text = f"{passed}/{total}" if total > 0 else "—"
            self._add_textbox(
                slide, Inches(5.5), y,
                Inches(1.5), Inches(0.45),
                text=score_text, font_size=13,
                color=Colors.SECONDARY_TEXT, alignment=PP_ALIGN.CENTER,
            )

            # Status label
            self._add_textbox(
                slide, Inches(7), y,
                Inches(5), Inches(0.45),
                text=style["label"], font_size=11,
                color=style["color"], bold=True,
            )

            y += Inches(0.6)

        # Recommendations
        if recommendations:
            y += Inches(0.3)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(12), Inches(0.35),
                text="Top Recommendations:", font_size=12, bold=True,
                color=Colors.TEAL,
            )
            y += Inches(0.4)
            for rec in recommendations[:4]:
                self._add_textbox(
                    slide, Inches(0.7), y,
                    Inches(11), Inches(0.3),
                    text=f"→ {rec}", font_size=10,
                )
                y += Inches(0.3)

        # Reference
        self._add_textbox(
            slide, self.MARGIN, Inches(6.5),
            Inches(12), Inches(0.35),
            text="Validated against Oracle Well-Architected Framework — docs.oracle.com/en/solutions/oci-best-practices/",
            font_size=10, italic=True, color=Colors.SECONDARY_TEXT,
        )

    def add_next_steps_slide(self, steps: list, contact_info: str = ""):
        """Slide 12: Next Steps."""
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Next Steps", margin=self.MARGIN)

        y = Inches(1.5)
        for i, step in enumerate(steps):
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                self.MARGIN, y,
                Inches(0.4), Inches(0.4),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = Colors.TEAL
            circle.line.fill.background()
            circle.text_frame.paragraphs[0].text = str(i + 1)
            circle.text_frame.paragraphs[0].font.size = Pt(12)
            circle.text_frame.paragraphs[0].font.bold = True
            circle.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            circle.text_frame.paragraphs[0].font.name = self.FONT
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Step text
            self._add_textbox(
                slide, Inches(1.1), y,
                Inches(10), Inches(0.4),
                text=step, font_size=13,
            )
            y += Inches(0.65)

        if contact_info:
            self._add_textbox(
                slide, self.MARGIN, Inches(6),
                Inches(12), Inches(0.4),
                text=contact_info, font_size=11,
                color=Colors.SECONDARY_TEXT,
            )

    # ---- Build from YAML spec ----

    @classmethod
    def from_spec(cls, spec: dict, template: Optional[str] = None) -> "OCIDeckGenerator":
        """Build a complete deck from a YAML specification."""
        # The MCP server passes a flat payload shape (customer_name, services,
        # cost_summary, ...) that doesn't match the proposal-spec structure this
        # method was originally written against. Adapt it so the deck renders
        # with actual content instead of only blank title/closing slides.
        if _is_flat_spec(spec):
            spec = _adapt_flat_spec(spec)

        meta = spec.get("metadata", {})
        gen = cls(
            customer=pick(meta, "customer", "customer_name"),
            project=pick(meta, "project", "project_name"),
            architect=pick(meta, "architect", "prepared_by"),
            firm=meta.get("firm", ""),
            template=template,
        )

        # Slide 1: Title
        gen.add_title_slide(subtitle=pick(meta, "subtitle", "title"))

        # Slide 2: Summary
        if "summary" in spec:
            s = spec["summary"]
            gen.add_summary_slide(
                why=pick(s, "why"),
                current_state=coerce_list(s.get("current_state")),
                target_state=pick(s, "target_state"),
                timeline=pick(s, "timeline"),
            )

        # Slide 3: Service Tiering (ECAL)
        if "service_tiering" in spec:
            gen.add_service_tiering_slide(spec["service_tiering"])

        # Slide 4: Architecture Principles (ECAL)
        if "architecture_principles" in spec:
            gen.add_architecture_principles_slide(spec["architecture_principles"])

        # Slide 5: Architecture
        if "architecture" in spec:
            a = spec["architecture"]
            gen.add_architecture_slide(
                diagram_path=a.get("diagram_path"),
                description=a.get("description", ""),
                visual=a.get("visual"),
            )

        # Slide 6: Decisions
        if "decisions" in spec:
            gen.add_decisions_slide(spec["decisions"])

        # Slide 7: HA/DR
        if "ha_dr" in spec:
            h = spec["ha_dr"]
            gen.add_ha_dr_slide(
                tiers=h.get("tiers", []),
                description=h.get("description", ""),
            )

        # Slide 8: Security
        if "security" in spec:
            s = spec["security"]
            gen.add_security_slide(
                controls=s.get("controls", {}),
                compliance=s.get("compliance", []),
            )

        # Slide 9: Environment Catalogue (ECAL)
        if "environment_catalogue" in spec:
            ec = spec["environment_catalogue"]
            gen.add_environment_catalogue_slide(
                environments=ec.get("environments", []),
                cost_notes=ec.get("cost_notes"),
            )

        # Slide 10: Cost
        if "cost" in spec:
            c = spec["cost"]
            gen.add_cost_slide(
                line_items=pick_list(c, "line_items"),
                assumptions=pick_list(c, "assumptions"),
                show_byol=c.get("show_byol", True),
            )

        # Slide 11: Cost Comparison (optional)
        if "cost_comparison" in spec:
            cc = spec["cost_comparison"]
            gen.add_cost_comparison_slide(
                rows=cc.get("rows", []),
                title=cc.get("title", "Cost Comparison"),
                col_headers=cc.get("col_headers"),
            )

        # Slide 12: Migration
        if "migration" in spec:
            m = spec["migration"]
            gen.add_migration_slide(
                phases=pick_list(m, "phases"),
                tools=pick_list(m, "tools"),
                downtime=pick(m, "downtime"),
            )

        # Slide 13: Operational RACI (ECAL)
        raci_payload = None
        for key in ("operational_raci", "raci", "operations_raci"):
            if key in spec and spec[key] is not None:
                raci_payload = spec[key]
                break
        if raci_payload is not None:
            model = "co_managed"
            raci_items = []
            if isinstance(raci_payload, str):
                model = raci_payload
            elif isinstance(raci_payload, list):
                raci_items = [r for r in raci_payload if isinstance(r, dict)]
            elif isinstance(raci_payload, dict):
                model = raci_payload.get("model") or raci_payload.get("engagement_model") or "co_managed"
                items = (
                    raci_payload.get("raci_items")
                    or raci_payload.get("items")
                    or raci_payload.get("activities")
                    or []
                )
                raci_items = [r for r in items if isinstance(r, dict)]
            if not raci_items:
                raci_items = _default_operational_raci(model)
            gen.add_operational_raci_slide(
                raci_items=raci_items,
                model=model,
            )

        # Slide 14: Risks
        if "risks" in spec:
            gen.add_risk_slide(spec["risks"])

        # Slide 15: WA Scorecard
        if "scorecard" in spec:
            sc = spec["scorecard"]
            gen.add_scorecard_slide(
                pillars=pick_list(sc, "pillars"),
                recommendations=pick_list(sc, "recommendations"),
            )

        # Slide 16: Next Steps
        if "next_steps" in spec:
            ns = spec["next_steps"]
            gen.add_next_steps_slide(
                steps=pick_list(ns, "steps", "next_steps"),
                contact_info=pick(ns, "contact_info"),
            )

        return gen


# ============================================================
# CLI
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate OCI architecture proposal slide deck (.pptx)"
    )
    parser.add_argument(
        "--spec", required=True,
        help="Path to YAML spec file with proposal data",
    )
    parser.add_argument(
        "--output", default="architecture-proposal.pptx",
        help="Output .pptx file path",
    )
    args = parser.parse_args()

    with open(args.spec, 'r') as f:
        spec = yaml.safe_load(f)

    gen = OCIDeckGenerator.from_spec(spec)
    gen.save(args.output)

    print(f"Generated: {args.output}")
    print(f"  Slides: {gen.slide_count}")
    print(f"  Template: Oracle FY26")
    customer = spec.get("metadata", {}).get("customer", "")
    if customer:
        print(f"  Customer: {customer}")


if __name__ == "__main__":
    main()
