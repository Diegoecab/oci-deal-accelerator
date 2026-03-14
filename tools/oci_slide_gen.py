#!/usr/bin/env python3
"""
OCI Deal Accelerator — Slide Deck Generator (.pptx)

Produces a 10-12 slide architecture proposal deck using OCI brand colors
and clean technical styling. No generic corporate templates.

Usage:
    python oci_slide_gen.py --spec proposal-data.yaml --output proposal.pptx

Or import and use programmatically:
    from oci_slide_gen import OCISlideGenerator
    gen = OCISlideGenerator(customer="Acme Corp", project="DB Migration")
    gen.add_summary_slide(...)
    gen.save("proposal.pptx")
"""

import yaml
import argparse
from datetime import datetime
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# OCI Brand Colors
# ============================================================

class Colors:
    # Primary
    DARK_BG = RGBColor(0x31, 0x2D, 0x2A)       # #312D2A — dark charcoal
    PRIMARY_TEXT = RGBColor(0x31, 0x2D, 0x2A)    # #312D2A — never pure black
    SECONDARY_TEXT = RGBColor(0x70, 0x66, 0x5E)  # #70665E — muted
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # OCI Accent Colors
    TEAL = RGBColor(0x2D, 0x59, 0x67)            # #2D5967 — primary accent
    COPPER = RGBColor(0xAA, 0x64, 0x3B)          # #AA643B — database
    PURPLE = RGBColor(0x80, 0x49, 0x98)          # #804998 — integration
    BURNT_ORANGE = RGBColor(0xAE, 0x56, 0x2C)    # #AE562C — warning/VCN

    # Status
    SUCCESS = RGBColor(0x5E, 0x96, 0x24)         # #5E9624 — muted green
    WARNING = RGBColor(0xAE, 0x56, 0x2C)         # #AE562C — burnt orange
    ERROR = RGBColor(0xC7, 0x46, 0x34)           # #C74634 — Oracle red

    # Table
    TABLE_ALT_ROW = RGBColor(0xF5, 0xF4, 0xF2)  # #F5F4F2 — light gray
    TABLE_HEADER_BG = RGBColor(0x2D, 0x59, 0x67) # teal


# ============================================================
# Slide Generator
# ============================================================

class OCISlideGenerator:
    """Generate OCI-branded architecture proposal slide decks."""

    SLIDE_WIDTH = Inches(13.333)   # Widescreen 16:9
    SLIDE_HEIGHT = Inches(7.5)
    MARGIN = Inches(0.5)
    FONT = "Segoe UI"

    def __init__(self, customer: str = "", project: str = "",
                 architect: str = "", firm: str = ""):
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        self.customer = customer
        self.project = project
        self.architect = architect
        self.firm = firm

    # ---- Helpers ----

    def _add_blank_slide(self):
        """Add a blank slide."""
        layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(layout)

    def _add_textbox(self, slide, left, top, width, height,
                     text="", font_size=12, bold=False, italic=False,
                     color=None, alignment=PP_ALIGN.LEFT,
                     font_name=None):
        """Add a text box to a slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font_name or self.FONT
        p.font.color.rgb = color or Colors.PRIMARY_TEXT
        p.alignment = alignment
        return txBox

    def _add_paragraph(self, text_frame, text="", font_size=12,
                       bold=False, italic=False, color=None,
                       alignment=PP_ALIGN.LEFT, space_before=0,
                       space_after=0, bullet=False):
        """Add a paragraph to an existing text frame."""
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = self.FONT
        p.font.color.rgb = color or Colors.PRIMARY_TEXT
        p.alignment = alignment
        if space_before:
            p.space_before = Pt(space_before)
        if space_after:
            p.space_after = Pt(space_after)
        if bullet:
            p.level = 0
        return p

    def _add_rect(self, slide, left, top, width, height,
                  fill_color=None, border_color=None):
        """Add a rectangle shape."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if border_color:
            shape.line.color.rgb = border_color
        else:
            shape.line.fill.background()
        return shape

    def _add_table(self, slide, rows, cols, left, top, width, height):
        """Add a table to a slide."""
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        return table

    def _style_table_cell(self, cell, text, font_size=10, bold=False,
                          color=None, bg_color=None, alignment=PP_ALIGN.LEFT):
        """Style a table cell."""
        cell.text = text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
            paragraph.font.name = self.FONT
            paragraph.font.color.rgb = color or Colors.PRIMARY_TEXT
            paragraph.alignment = alignment
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if bg_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color

    def _add_title_bar(self, slide, title_text):
        """Add a consistent title bar at top of content slides."""
        # Title text
        self._add_textbox(
            slide, self.MARGIN, Inches(0.3),
            Inches(10), Inches(0.6),
            text=title_text, font_size=24, bold=True,
            color=Colors.PRIMARY_TEXT,
        )
        # Accent line under title
        self._add_rect(
            slide, self.MARGIN, Inches(0.85),
            Inches(2), Inches(0.04),
            fill_color=Colors.TEAL,
        )

    # ---- Slide Methods ----

    def add_title_slide(self, subtitle: str = ""):
        """Slide 1: Title slide with dark background."""
        slide = self._add_blank_slide()

        # Dark background
        self._add_rect(
            slide, Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            fill_color=Colors.DARK_BG,
        )

        # Teal accent bar
        self._add_rect(
            slide, Inches(0.8), Inches(2.2),
            Inches(0.08), Inches(2.5),
            fill_color=Colors.TEAL,
        )

        # Customer + Project
        title = f"{self.customer}" if self.customer else "Architecture Proposal"
        self._add_textbox(
            slide, Inches(1.2), Inches(2.2),
            Inches(10), Inches(0.7),
            text=title, font_size=32, bold=True,
            color=Colors.WHITE,
        )

        # Project name
        if self.project:
            self._add_textbox(
                slide, Inches(1.2), Inches(2.9),
                Inches(10), Inches(0.6),
                text=self.project, font_size=22,
                color=Colors.WHITE,
            )

        # Subtitle
        date_str = datetime.now().strftime("%B %Y")
        sub = subtitle or f"Architecture Proposal — {date_str}"
        self._add_textbox(
            slide, Inches(1.2), Inches(3.6),
            Inches(10), Inches(0.5),
            text=sub, font_size=16,
            color=RGBColor(0x9E, 0x98, 0x92),  # muted gray
        )

        # Architect / Firm
        if self.architect or self.firm:
            info = f"{self.architect}"
            if self.firm:
                info += f"  |  {self.firm}" if self.architect else self.firm
            self._add_textbox(
                slide, Inches(1.2), Inches(5.5),
                Inches(10), Inches(0.4),
                text=info, font_size=12,
                color=RGBColor(0x9E, 0x98, 0x92),
            )

        # "Prepared with OCI Deal Accelerator"
        self._add_textbox(
            slide, Inches(1.2), Inches(6.2),
            Inches(10), Inches(0.3),
            text="Prepared with OCI Deal Accelerator", font_size=9,
            italic=True, color=RGBColor(0x70, 0x66, 0x5E),
        )

    def add_summary_slide(self, why: str, current_state: list,
                          target_state: str, timeline: str):
        """Slide 2: Engagement Summary."""
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Engagement Summary")

        y = Inches(1.2)

        # Why
        box = self._add_textbox(
            slide, self.MARGIN, y, Inches(11), Inches(0.5),
            text=why, font_size=14, color=Colors.TEAL, bold=True,
        )
        y += Inches(0.7)

        # Current state
        self._add_textbox(
            slide, self.MARGIN, y, Inches(5), Inches(0.4),
            text="Current State", font_size=14, bold=True,
        )
        y += Inches(0.45)

        for item in current_state:
            self._add_textbox(
                slide, Inches(0.7), y, Inches(10), Inches(0.35),
                text=f"•  {item}", font_size=12,
            )
            y += Inches(0.35)

        y += Inches(0.3)

        # Target state
        self._add_textbox(
            slide, self.MARGIN, y, Inches(5), Inches(0.4),
            text="Target State", font_size=14, bold=True,
        )
        y += Inches(0.45)
        self._add_textbox(
            slide, Inches(0.7), y, Inches(10), Inches(0.4),
            text=target_state, font_size=12, color=Colors.TEAL,
        )
        y += Inches(0.5)

        # Timeline
        self._add_textbox(
            slide, self.MARGIN, y, Inches(10), Inches(0.4),
            text=f"Timeline: {timeline}", font_size=12, bold=True,
        )

    def add_architecture_slide(self, diagram_path: Optional[str] = None,
                                description: str = ""):
        """Slide 3: Architecture Overview with diagram."""
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Architecture Overview")

        if diagram_path:
            try:
                slide.shapes.add_picture(
                    diagram_path,
                    Inches(0.5), Inches(1.1),
                    Inches(12.3), Inches(6),
                )
            except Exception:
                self._add_textbox(
                    slide, Inches(1), Inches(2.5),
                    Inches(10), Inches(2),
                    text=f"[Diagram: {diagram_path}]\n\n{description}",
                    font_size=14, color=Colors.SECONDARY_TEXT,
                )
        else:
            self._add_textbox(
                slide, Inches(1), Inches(2),
                Inches(11), Inches(4),
                text=description or "[Insert architecture diagram — export from .drawio file]",
                font_size=14, color=Colors.SECONDARY_TEXT, italic=True,
                alignment=PP_ALIGN.CENTER,
            )

    def add_decisions_slide(self, decisions: list):
        """Slide 4: Architecture Decisions table.

        decisions: list of {"decision": str, "rationale": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Architecture Decisions")

        rows = len(decisions) + 1
        table = self._add_table(
            slide, rows, 2,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.5 * rows),
        )

        # Set column widths
        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(7.5)

        # Header row
        self._style_table_cell(
            table.cell(0, 0), "Decision", font_size=11, bold=True,
            color=Colors.WHITE, bg_color=Colors.TEAL,
        )
        self._style_table_cell(
            table.cell(0, 1), "Rationale", font_size=11, bold=True,
            color=Colors.WHITE, bg_color=Colors.TEAL,
        )

        # Data rows
        for i, dec in enumerate(decisions):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(
                table.cell(row_idx, 0), dec["decision"],
                font_size=11, bold=True, bg_color=bg,
            )
            self._style_table_cell(
                table.cell(row_idx, 1), dec["rationale"],
                font_size=10, bg_color=bg,
            )

    def add_ha_dr_slide(self, tiers: list, description: str = ""):
        """Slide 5: HA/DR table.

        tiers: list of {"tier": str, "technology": str, "rto": str, "rpo": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "High Availability & Disaster Recovery")

        if description:
            self._add_textbox(
                slide, self.MARGIN, Inches(1.2),
                Inches(12), Inches(0.5),
                text=description, font_size=12,
            )

        rows = len(tiers) + 1
        table = self._add_table(
            slide, rows, 4,
            self.MARGIN, Inches(1.9),
            Inches(12), Inches(0.45 * rows),
        )

        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(4.5)
        table.columns[2].width = Inches(2.5)
        table.columns[3].width = Inches(2.5)

        headers = ["Tier", "Technology", "RTO", "RPO"]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER,
            )

        for i, tier in enumerate(tiers):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(table.cell(row_idx, 0), tier["tier"], font_size=10, bold=True, bg_color=bg)
            self._style_table_cell(table.cell(row_idx, 1), tier["technology"], font_size=10, bg_color=bg)
            self._style_table_cell(table.cell(row_idx, 2), tier["rto"], font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)
            self._style_table_cell(table.cell(row_idx, 3), tier["rpo"], font_size=10, bg_color=bg, alignment=PP_ALIGN.CENTER)

    def add_security_slide(self, controls: dict, compliance: list = None):
        """Slide 6: Security & Compliance.

        controls: {"identity": [...], "network": [...], "database": [...], "monitoring": [...]}
        compliance: ["PCI-DSS", "SOC2", ...] or None
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Security & Compliance")

        y = Inches(1.2)

        # Compliance badges
        if compliance:
            x = self.MARGIN
            for framework in compliance:
                shape = self._add_rect(
                    slide, x, y, Inches(1.5), Inches(0.4),
                    fill_color=Colors.SUCCESS,
                )
                shape.text_frame.paragraphs[0].text = f"✓ {framework}"
                shape.text_frame.paragraphs[0].font.size = Pt(10)
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
                shape.text_frame.paragraphs[0].font.name = self.FONT
                shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                x += Inches(1.7)
            y += Inches(0.7)

        # Controls in columns
        col_x = self.MARGIN
        col_width = Inches(3)
        for area, items in controls.items():
            area_title = area.replace("_", " ").title()
            self._add_textbox(
                slide, col_x, y, col_width, Inches(0.35),
                text=area_title, font_size=12, bold=True,
                color=Colors.TEAL,
            )
            item_y = y + Inches(0.4)
            for item in items:
                self._add_textbox(
                    slide, col_x + Inches(0.1), item_y,
                    col_width - Inches(0.1), Inches(0.3),
                    text=f"• {item}", font_size=10,
                )
                item_y += Inches(0.3)
            col_x += Inches(3.2)

    def add_cost_slide(self, line_items: list, assumptions: list = None,
                       show_byol: bool = True):
        """Slide 7: Cost Estimate.

        line_items: list of {"component": str, "monthly_payg": str,
                            "monthly_byol": str (optional), "notes": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Cost Estimate")

        cols = 4 if show_byol else 3
        rows = len(line_items) + 1
        table = self._add_table(
            slide, rows, cols,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.4 * rows),
        )

        if show_byol:
            table.columns[0].width = Inches(3.5)
            table.columns[1].width = Inches(2.5)
            table.columns[2].width = Inches(2.5)
            table.columns[3].width = Inches(3.5)
            headers = ["Component", "Monthly (PAYG)", "Monthly (BYOL)", "Notes"]
        else:
            table.columns[0].width = Inches(4)
            table.columns[1].width = Inches(3)
            table.columns[2].width = Inches(5)
            headers = ["Component", "Monthly", "Notes"]

        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
            )

        for i, item in enumerate(line_items):
            row_idx = i + 1
            is_total = "total" in item.get("component", "").lower()
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 and not is_total else None
            if is_total:
                bg = Colors.TEAL

            self._style_table_cell(
                table.cell(row_idx, 0), item["component"],
                font_size=10, bold=is_total,
                color=Colors.WHITE if is_total else None,
                bg_color=bg,
            )
            self._style_table_cell(
                table.cell(row_idx, 1), item.get("monthly_payg", ""),
                font_size=10, bold=is_total,
                color=Colors.WHITE if is_total else None,
                bg_color=bg, alignment=PP_ALIGN.RIGHT,
            )
            if show_byol:
                self._style_table_cell(
                    table.cell(row_idx, 2), item.get("monthly_byol", ""),
                    font_size=10, bold=is_total,
                    color=Colors.WHITE if is_total else None,
                    bg_color=bg, alignment=PP_ALIGN.RIGHT,
                )
                self._style_table_cell(
                    table.cell(row_idx, 3), item.get("notes", ""),
                    font_size=9, bg_color=bg,
                    color=Colors.WHITE if is_total else Colors.SECONDARY_TEXT,
                )
            else:
                self._style_table_cell(
                    table.cell(row_idx, 2), item.get("notes", ""),
                    font_size=9, bg_color=bg,
                    color=Colors.WHITE if is_total else Colors.SECONDARY_TEXT,
                )

        # Assumptions
        if assumptions:
            table_bottom = Inches(1.2) + Inches(0.4 * rows) + Inches(0.3)
            self._add_textbox(
                slide, self.MARGIN, table_bottom,
                Inches(12), Inches(0.3),
                text="Assumptions:", font_size=9, bold=True,
                color=Colors.SECONDARY_TEXT, italic=True,
            )
            for idx, assumption in enumerate(assumptions):
                self._add_textbox(
                    slide, Inches(0.7), table_bottom + Inches(0.3 + idx * 0.25),
                    Inches(11), Inches(0.25),
                    text=f"• {assumption}", font_size=8, italic=True,
                    color=Colors.SECONDARY_TEXT,
                )

    def add_migration_slide(self, phases: list, tools: list = None,
                            downtime: str = ""):
        """Slide 9: Migration Approach.

        phases: list of {"name": str, "weeks": str, "tasks": [...]}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Migration Approach")

        y = Inches(1.3)
        phase_colors = [Colors.TEAL, Colors.COPPER, Colors.PURPLE, Colors.BURNT_ORANGE]

        # Phase timeline bars
        bar_left = Inches(2.5)
        bar_width = Inches(9.5)
        for i, phase in enumerate(phases):
            color = phase_colors[i % len(phase_colors)]

            # Phase label
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(2), Inches(0.45),
                text=phase["name"], font_size=11, bold=True,
            )

            # Phase bar
            bar = self._add_rect(
                slide, bar_left, y + Inches(0.05),
                bar_width * 0.9, Inches(0.35),
                fill_color=color,
            )
            bar.text_frame.paragraphs[0].text = phase.get("weeks", "")
            bar.text_frame.paragraphs[0].font.size = Pt(9)
            bar.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            bar.text_frame.paragraphs[0].font.name = self.FONT
            bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Tasks below bar
            if phase.get("tasks"):
                tasks_text = " | ".join(phase["tasks"][:4])
                self._add_textbox(
                    slide, bar_left, y + Inches(0.4),
                    bar_width, Inches(0.3),
                    text=tasks_text, font_size=8, italic=True,
                    color=Colors.SECONDARY_TEXT,
                )

            y += Inches(0.9)

        # Tools
        if tools:
            y += Inches(0.2)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(12), Inches(0.35),
                text=f"Migration Tools: {', '.join(tools)}",
                font_size=11, bold=True, color=Colors.TEAL,
            )

        # Downtime approach
        if downtime:
            y += Inches(0.5)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(12), Inches(0.35),
                text=f"Downtime Approach: {downtime}",
                font_size=11, color=Colors.PRIMARY_TEXT,
            )

    def add_risk_slide(self, risks: list):
        """Slide 10: Risk Register.

        risks: list of {"risk": str, "severity": "HIGH"|"MEDIUM"|"LOW", "mitigation": str}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Risk Register")

        rows = len(risks) + 1
        table = self._add_table(
            slide, rows, 3,
            self.MARGIN, Inches(1.2),
            Inches(12), Inches(0.5 * rows),
        )

        table.columns[0].width = Inches(4.5)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(6)

        headers = ["Risk", "Severity", "Mitigation"]
        for j, h in enumerate(headers):
            self._style_table_cell(
                table.cell(0, j), h, font_size=11, bold=True,
                color=Colors.WHITE, bg_color=Colors.TEAL,
                alignment=PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT,
            )

        severity_colors = {
            "HIGH": Colors.ERROR,
            "MEDIUM": Colors.WARNING,
            "LOW": Colors.SUCCESS,
        }

        for i, risk in enumerate(risks):
            row_idx = i + 1
            bg = Colors.TABLE_ALT_ROW if row_idx % 2 == 0 else None
            self._style_table_cell(
                table.cell(row_idx, 0), risk["risk"],
                font_size=10, bg_color=bg,
            )
            sev = risk.get("severity", "MEDIUM").upper()
            self._style_table_cell(
                table.cell(row_idx, 1), sev,
                font_size=10, bold=True,
                color=severity_colors.get(sev, Colors.WARNING),
                bg_color=bg, alignment=PP_ALIGN.CENTER,
            )
            self._style_table_cell(
                table.cell(row_idx, 2), risk["mitigation"],
                font_size=10, bg_color=bg,
            )

    def add_scorecard_slide(self, pillars: list, recommendations: list = None):
        """Slide 11: Well-Architected Scorecard.

        pillars: list of {"name": str, "status": "PASS"|"PASS_WITH_RECOMMENDATIONS"|"GAPS_IDENTIFIED"|"NOT_APPLICABLE",
                         "passed": int, "total": int}
        """
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Well-Architected Scorecard")

        status_styles = {
            "PASS": {"icon": "✓", "color": Colors.SUCCESS, "label": "PASS"},
            "PASS_WITH_RECOMMENDATIONS": {"icon": "⚠", "color": Colors.WARNING, "label": "PASS WITH RECOMMENDATIONS"},
            "GAPS_IDENTIFIED": {"icon": "✗", "color": Colors.ERROR, "label": "GAPS IDENTIFIED"},
            "NOT_APPLICABLE": {"icon": "—", "color": Colors.SECONDARY_TEXT, "label": "N/A"},
        }

        y = Inches(1.4)
        for pillar in pillars:
            style = status_styles.get(pillar["status"], status_styles["PASS"])

            # Status indicator pill
            pill = self._add_rect(
                slide, self.MARGIN, y,
                Inches(0.5), Inches(0.45),
                fill_color=style["color"],
            )
            pill.text_frame.paragraphs[0].text = style["icon"]
            pill.text_frame.paragraphs[0].font.size = Pt(14)
            pill.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            pill.text_frame.paragraphs[0].font.name = self.FONT
            pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Pillar name
            self._add_textbox(
                slide, Inches(1.2), y,
                Inches(4), Inches(0.45),
                text=pillar["name"], font_size=13, bold=True,
            )

            # Score
            score_text = f"{pillar['passed']}/{pillar['total']}" if pillar['total'] > 0 else "—"
            self._add_textbox(
                slide, Inches(5.5), y,
                Inches(1.5), Inches(0.45),
                text=score_text, font_size=13,
                color=Colors.SECONDARY_TEXT, alignment=PP_ALIGN.CENTER,
            )

            # Status label
            self._add_textbox(
                slide, Inches(7), y,
                Inches(5), Inches(0.45),
                text=style["label"], font_size=11,
                color=style["color"], bold=True,
            )

            y += Inches(0.6)

        # Recommendations
        if recommendations:
            y += Inches(0.3)
            self._add_textbox(
                slide, self.MARGIN, y,
                Inches(12), Inches(0.35),
                text="Top Recommendations:", font_size=12, bold=True,
                color=Colors.TEAL,
            )
            y += Inches(0.4)
            for rec in recommendations[:4]:
                self._add_textbox(
                    slide, Inches(0.7), y,
                    Inches(11), Inches(0.3),
                    text=f"→ {rec}", font_size=10,
                )
                y += Inches(0.3)

        # Reference
        self._add_textbox(
            slide, self.MARGIN, Inches(6.5),
            Inches(12), Inches(0.3),
            text="Validated against Oracle Well-Architected Framework — docs.oracle.com/en/solutions/oci-best-practices/",
            font_size=8, italic=True, color=Colors.SECONDARY_TEXT,
        )

    def add_next_steps_slide(self, steps: list, contact_info: str = ""):
        """Slide 12: Next Steps."""
        slide = self._add_blank_slide()
        self._add_title_bar(slide, "Next Steps")

        y = Inches(1.5)
        for i, step in enumerate(steps):
            # Number circle
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                self.MARGIN, y,
                Inches(0.4), Inches(0.4),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = Colors.TEAL
            circle.line.fill.background()
            circle.text_frame.paragraphs[0].text = str(i + 1)
            circle.text_frame.paragraphs[0].font.size = Pt(12)
            circle.text_frame.paragraphs[0].font.bold = True
            circle.text_frame.paragraphs[0].font.color.rgb = Colors.WHITE
            circle.text_frame.paragraphs[0].font.name = self.FONT
            circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Step text
            self._add_textbox(
                slide, Inches(1.1), y,
                Inches(10), Inches(0.4),
                text=step, font_size=13,
            )
            y += Inches(0.65)

        if contact_info:
            self._add_textbox(
                slide, self.MARGIN, Inches(6),
                Inches(12), Inches(0.4),
                text=contact_info, font_size=11,
                color=Colors.SECONDARY_TEXT,
            )

    # ---- Build from YAML spec ----

    @classmethod
    def from_spec(cls, spec: dict) -> "OCISlideGenerator":
        """Build a complete deck from a YAML specification."""
        meta = spec.get("metadata", {})
        gen = cls(
            customer=meta.get("customer", ""),
            project=meta.get("project", ""),
            architect=meta.get("architect", ""),
            firm=meta.get("firm", ""),
        )

        # Slide 1: Title
        gen.add_title_slide(subtitle=meta.get("subtitle", ""))

        # Slide 2: Summary
        if "summary" in spec:
            s = spec["summary"]
            gen.add_summary_slide(
                why=s.get("why", ""),
                current_state=s.get("current_state", []),
                target_state=s.get("target_state", ""),
                timeline=s.get("timeline", ""),
            )

        # Slide 3: Architecture
        if "architecture" in spec:
            a = spec["architecture"]
            gen.add_architecture_slide(
                diagram_path=a.get("diagram_path"),
                description=a.get("description", ""),
            )

        # Slide 4: Decisions
        if "decisions" in spec:
            gen.add_decisions_slide(spec["decisions"])

        # Slide 5: HA/DR
        if "ha_dr" in spec:
            h = spec["ha_dr"]
            gen.add_ha_dr_slide(
                tiers=h.get("tiers", []),
                description=h.get("description", ""),
            )

        # Slide 6: Security
        if "security" in spec:
            s = spec["security"]
            gen.add_security_slide(
                controls=s.get("controls", {}),
                compliance=s.get("compliance", []),
            )

        # Slide 7: Cost
        if "cost" in spec:
            c = spec["cost"]
            gen.add_cost_slide(
                line_items=c.get("line_items", []),
                assumptions=c.get("assumptions", []),
                show_byol=c.get("show_byol", True),
            )

        # Slide 9: Migration
        if "migration" in spec:
            m = spec["migration"]
            gen.add_migration_slide(
                phases=m.get("phases", []),
                tools=m.get("tools", []),
                downtime=m.get("downtime", ""),
            )

        # Slide 10: Risks
        if "risks" in spec:
            gen.add_risk_slide(spec["risks"])

        # Slide 11: WA Scorecard
        if "scorecard" in spec:
            sc = spec["scorecard"]
            gen.add_scorecard_slide(
                pillars=sc.get("pillars", []),
                recommendations=sc.get("recommendations", []),
            )

        # Slide 12: Next Steps
        if "next_steps" in spec:
            ns = spec["next_steps"]
            gen.add_next_steps_slide(
                steps=ns.get("steps", []),
                contact_info=ns.get("contact_info", ""),
            )

        return gen

    # ---- Save ----

    def save(self, filepath: str):
        """Save the presentation to a .pptx file."""
        self.prs.save(filepath)

    @property
    def slide_count(self):
        return len(self.prs.slides)


# ============================================================
# CLI
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate OCI architecture proposal slide deck (.pptx)"
    )
    parser.add_argument(
        "--spec", required=True,
        help="Path to YAML spec file with proposal data",
    )
    parser.add_argument(
        "--output", default="architecture-proposal.pptx",
        help="Output .pptx file path",
    )
    args = parser.parse_args()

    with open(args.spec, 'r') as f:
        spec = yaml.safe_load(f)

    gen = OCISlideGenerator.from_spec(spec)
    gen.save(args.output)

    print(f"Generated: {args.output}")
    print(f"  Slides: {gen.slide_count}")
    customer = spec.get("metadata", {}).get("customer", "")
    if customer:
        print(f"  Customer: {customer}")


if __name__ == "__main__":
    main()
