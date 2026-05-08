import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


PROJECT_ROOT = Path(__file__).resolve().parent.parent
TOOLS_DIR = PROJECT_ROOT / "tools"

if str(TOOLS_DIR) not in sys.path:
    sys.path.insert(0, str(TOOLS_DIR))

from pptx_validator import validate_pptx  # noqa: E402


def test_validate_pptx_reads_text_boxes_tables_and_flags_forbidden(tmp_path):
    pptx = tmp_path / "validation-fixture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.4)).text = "TCO Crossover"
    table = slide.shapes.add_table(1, 1, Inches(0.5), Inches(1.0), Inches(4), Inches(0.5)).table
    table.cell(0, 0).text = "Storage break-even"
    slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4), Inches(0.4)).text = "OCI Annual"
    last = prs.slides.add_slide(prs.slide_layouts[6])
    last.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.4)).text = "Commercial Disclaimer"
    prs.save(pptx)

    report = validate_pptx(
        pptx,
        expected_titles=["TCO Crossover"],
        required_phrases=["Storage break-even"],
        forbidden_phrases=["OCI Annual"],
        disclaimer_last=True,
    )

    assert report["status"] == "fail"
    assert any(issue["code"] == "FORBIDDEN_PHRASE" for issue in report["issues"])


def test_validate_pptx_passes_required_and_disclaimer_last(tmp_path):
    pptx = tmp_path / "validation-pass.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7), Inches(0.4)).text = (
        "BYOL/PAYG model Discount GoldenGate bridge duration "
        "Workload ECPU demand ECPU capacity Storage break-even Crossover"
    )
    last = prs.slides.add_slide(prs.slide_layouts[6])
    last.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.4)).text = "Commercial Disclaimer"
    prs.save(pptx)

    report = validate_pptx(
        pptx,
        required_phrases=["BYOL/PAYG model", "Storage break-even"],
        forbidden_phrases=["FTE-year"],
        disclaimer_last=True,
    )

    assert report["status"] == "pass"
