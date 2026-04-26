#!/usr/bin/env python3
"""
Generate reproducible random evaluation cases plus local artifacts.

This script exists to stress the skill with sparse-but-valid prompts:
- Creates a Spanish discovery prompt
- Derives minimal local specs/payloads for deck, diagram, business case, and BOM
- Renders local artifacts (deck, drawio, business case, BOM, AppCA BOM, PDF)
- Produces lightweight artifact quality metrics for later MCP comparison
"""

from __future__ import annotations

import argparse
import json
import random
import re
import sys
from pathlib import Path
from zipfile import ZipFile

import yaml
from openpyxl import load_workbook
from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parent.parent
TOOLS_DIR = PROJECT_ROOT / "tools"
if str(TOOLS_DIR) not in sys.path:
    sys.path.insert(0, str(TOOLS_DIR))

from oci_bizcase_gen import BusinessCaseDeckGenerator
from oci_bom_gen import OCIBomGenerator
from oci_deck_gen import OCIDeckGenerator, _enrich_partial_proposal_spec
from oci_diagram_gen import OCIDiagramGenerator
from oci_pdf_gen import OCIPDFGenerator


INDUSTRIES = [
    ("banca", "banking regulator", "PCI-DSS"),
    ("retail", "consumer peak events", "PCI-DSS"),
    ("salud", "patient data controls", "HIPAA-aligned controls"),
    ("manufactura", "plant downtime sensitivity", "ISO 27001"),
    ("logística", "cross-border operations", "SOC 2"),
]

TEMPLATES = [
    {
        "id": "adbs",
        "workload": "e-commerce platform",
        "current": "dos bases Oracle 19c on-prem y VMs para web",
        "target": "ADB-S para transaccional, compute flexible para app y Object Storage para backups",
        "services": [
            {"name": "ADB-S", "sku": "B95701", "kind": "database"},
            {"name": "ADB Storage", "sku": "B95706", "kind": "database_storage"},
            {"name": "ADB Backup", "sku": "B95754a", "kind": "database_storage"},
            {"name": "Block Volume", "sku": "B91961", "kind": "storage"},
            {"name": "FastConnect 1 Gbps", "sku": "B88325", "kind": "network"},
        ],
    },
    {
        "id": "exacs",
        "workload": "core banking database platform",
        "current": "Exadata X8M on-prem con base crítica y reporting separado",
        "target": "ExaCS X11M BYOL para base crítica, DR cross-region y Object Storage para respaldo",
        "services": [
            {"name": "Exadata Base System", "sku": "B90777", "kind": "database"},
            {"name": "Exadata DB Server X11M", "sku": "B110627", "kind": "database"},
            {"name": "Exadata Storage Server X11M", "sku": "B110629", "kind": "database"},
            {"name": "Exadata Dedicated ECPU BYOL", "sku": "B110632", "kind": "database"},
        ],
    },
    {
        "id": "analytics",
        "workload": "analytics and data science stack",
        "current": "ETL batch en Hadoop heredado y notebooks aislados",
        "target": "Big Data Service + Data Science + Data Flow con Object Storage como data lake",
        "services": [
            {"name": "Big Data Service Standard", "sku": "B91128", "kind": "analytics"},
            {"name": "Data Science Notebook Estimate", "sku": "EST-DS-NOTEBOOK", "kind": "analytics"},
            {"name": "Data Science Model Estimate", "sku": "EST-DS-MODEL", "kind": "analytics"},
            {"name": "Data Flow Spark Estimate", "sku": "EST-DF-SPARK", "kind": "analytics"},
            {"name": "Block Volume", "sku": "B91961", "kind": "storage"},
        ],
    },
]

PRIMARY_REGIONS = ["us-ashburn-1", "us-phoenix-1", "eu-frankfurt-1", "sa-saopaulo-1"]
DR_REGIONS = ["us-phoenix-1", "us-sanjose-1", "eu-amsterdam-1", "sa-vinhedo-1"]


def slugify(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")


def build_case(seed: int) -> dict:
    rng = random.Random(seed)
    industry, urgency, compliance = rng.choice(INDUSTRIES)
    template = rng.choice(TEMPLATES)
    customer = f"{rng.choice(['Nova', 'Andes', 'Pacific', 'Vector', 'Apex'])} {rng.choice(['Digital', 'Holdings', 'Retail', 'Bank', 'Health'])}"
    primary_region = rng.choice(PRIMARY_REGIONS)
    dr_region = rng.choice([r for r in DR_REGIONS if r != primary_region])
    timeline_weeks = rng.choice([12, 14, 16, 20, 24])
    peak = rng.choice(["2x", "3x", "4x"])
    budget = rng.choice(["ajustado", "moderado", "sujeto a aprobación trimestral"])
    team = rng.choice(["2 DBAs y 1 sysadmin", "1 DBA senior y equipo de aplicaciones", "equipo pequeño sin experiencia OCI"])
    quantities = {}

    for service in template["services"]:
        sku = service["sku"]
        if sku in {"B95701", "B110632"}:
            quantities[sku] = rng.choice([16, 24, 32, 64])
        elif sku in {"B95706", "B95754a", "B91961"}:
            quantities[sku] = rng.choice([512, 1024, 2048, 4096])
        elif sku in {"B90777", "B88325"}:
            quantities[sku] = 1
        elif sku in {"B110627"}:
            quantities[sku] = 2
        elif sku in {"B110629"}:
            quantities[sku] = 3
        else:
            quantities[sku] = rng.choice([2, 4, 8, 16])

    prompt = (
        f"Cliente: {customer}. Industria: {industry}. Driver principal: reducir costo y mejorar DR.\n"
        f"Estado actual: {template['current']}. Objetivo: {template['target']}.\n"
        f"Región primaria: {primary_region}. Región DR: {dr_region}. Ventana objetivo: {timeline_weeks} semanas.\n"
        f"Compliance: {compliance}. Sensibilidad de presupuesto: {budget}. Equipo actual: {team}.\n"
        f"Picos esperados: {peak}. Contexto adicional: {urgency}."
    )

    customer_id = slugify(customer)
    summary_current = [
        template["current"],
        f"Primary region target: {primary_region}",
        f"DR region target: {dr_region}",
        f"Compliance baseline: {compliance}",
        f"Peak growth assumption: {peak}",
    ]
    summary_target = (
        f"{template['target']}. Primary region {primary_region}; DR region {dr_region}. "
        f"Timeline target {timeline_weeks} weeks."
    )

    minimal_proposal_spec = {
        "metadata": {
            "customer": customer,
            "project": template["workload"].title(),
            "subtitle": f"OCI proposal for {template['workload']}",
        },
        "summary": {
            "why": "Modernize the platform while protecting business continuity and commercial efficiency.",
            "current_state": summary_current,
            "target_state": summary_target,
            "timeline": f"{timeline_weeks} weeks",
        },
    }

    diagram_services = []
    for idx, service in enumerate(template["services"][:3], 1):
        svc_type = "database" if service["kind"] in {"database", "database_storage"} else "compute"
        if service["kind"] == "network":
            svc_type = "fastconnect"
        elif service["kind"] == "storage":
            svc_type = "object_storage"
        diagram_services.append({
            "id": f"svc{idx}",
            "label": service["name"],
            "type": svc_type,
        })

    diagram_spec = {
        "title": f"{customer} — {template['workload']}",
        "external": [
            {"id": "users", "label": "Enterprise\nUsers", "icon": "user", "x": 30, "y": 260, "w": 80, "h": 80},
        ],
        "tenancy": {
            "label": f"OCI Tenancy — {customer}",
            "regions": [
                {
                    "id": "primary",
                    "label": f"Region — {primary_region} (Primary)",
                    "primary": True,
                    "vcns": [
                        {
                            "id": "vcn1",
                            "label": "Application VCN",
                            "subnets": [
                                {"id": "subnet1", "label": "Application / Data Subnet", "services": diagram_services},
                            ],
                        }
                    ],
                },
                {
                    "id": "dr",
                    "label": f"Region — {dr_region} (DR)",
                    "primary": False,
                    "vcns": [
                        {
                            "id": "vcn2",
                            "label": "DR VCN",
                            "subnets": [
                                {
                                    "id": "subnet2",
                                    "label": "Standby Subnet",
                                    "services": [
                                        {
                                            "id": "drsvc1",
                                            "label": "Standby / DR",
                                            "type": "database",
                                        }
                                    ],
                                },
                            ],
                        }
                    ],
                },
            ],
        },
        "connections": [
            {"from": "users", "to": "svc1", "type": "standard", "label": "Private access"},
            {"from": "svc1", "to": "drsvc1", "type": "data", "label": "Replication"},
        ],
    }

    business_case_spec = {
        "customer_name": customer,
        "executive_summary": prompt.replace("\n", " "),
    }

    services_payload = [
        {"sku": service["sku"], "quantity": quantities[service["sku"]]}
        for service in template["services"]
    ]

    bom_spec = {
        "bom": {
            "customer_name": customer,
            "project_name": template["workload"].title(),
            "prepared_by": "Codex Evaluation Harness",
            "currency": "USD",
            "line_items": [
                {"sku": service["sku"], "qty": quantities[service["sku"]], "months": 12, "discount": 0.0}
                for service in template["services"]
            ],
            "notes": [
                f"Generated from seed {seed}",
                "Commercial validation required before quoting",
            ],
        }
    }

    return {
        "seed": seed,
        "customer": customer,
        "customer_id": customer_id,
        "prompt": prompt,
        "proposal_spec": minimal_proposal_spec,
        "diagram_spec": diagram_spec,
        "business_case_spec": business_case_spec,
        "bom_spec": bom_spec,
        "mcp_payloads": {
            "deck": {
                "customer_id": customer_id,
                "preview": False,
                "tier": "standard",
                "spec": minimal_proposal_spec,
            },
            "diagram": {
                "customer_id": customer_id,
                "preview": False,
                "spec": diagram_spec,
            },
            "business_case": {
                "customer_id": customer_id,
                "preview": False,
                "discovery_notes": prompt,
            },
            "bom": {
                "customer_id": customer_id,
                "preview": False,
                "currency": "USD",
                "discount_pct": 0.0,
                "services": services_payload,
            },
            "bom_appca": {
                "customer_id": customer_id,
                "preview": False,
                "currency": "USD",
                "discount_pct": 0.0,
                "services": services_payload,
            },
        },
    }


def save_yaml(path: Path, payload: dict):
    path.write_text(yaml.safe_dump(payload, sort_keys=False), encoding="utf-8")


def analyze_pptx(path: Path) -> dict:
    prs = Presentation(path)
    slides = []
    blank = 0
    for slide in prs.slides:
        texts = []
        tables = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                tables += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            texts.append(cell.text.strip())
        if len(texts) <= 1 and tables == 0:
            blank += 1
        slides.append({"text_items": len(texts), "tables": tables, "sample": texts[:6]})
    return {"slide_count": len(prs.slides), "blank_slides": blank, "slides": slides}


def analyze_xlsx(path: Path) -> dict:
    wb = load_workbook(path, data_only=False)
    sheets = []
    empty = 0
    for ws in wb.worksheets:
        nonempty = 0
        for row in ws.iter_rows():
            for cell in row:
                if cell.value not in (None, ""):
                    nonempty += 1
        if nonempty == 0:
            empty += 1
        sheets.append({"name": ws.title, "rows": ws.max_row, "cols": ws.max_column, "nonempty_cells": nonempty})
    return {"sheet_count": len(wb.worksheets), "empty_sheets": empty, "sheets": sheets}


def analyze_drawio(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    service_labels = len(re.findall(r'value="[^"]+"', text))
    return {"bytes": path.stat().st_size, "cell_values": service_labels}


def analyze_pdf(path: Path) -> dict:
    raw = path.read_bytes()
    page_count = raw.count(b"/Type /Page")
    snippets = re.findall(rb"[A-Za-z][A-Za-z0-9 ,.:()/-]{8,}", raw)
    decoded = []
    for snippet in snippets[:20]:
        try:
            decoded.append(snippet.decode("utf-8"))
        except UnicodeDecodeError:
            decoded.append(snippet.decode("latin-1", errors="ignore"))
    return {"bytes": path.stat().st_size, "page_count_estimate": page_count, "snippets": decoded[:10]}


def render_local(case: dict, out_dir: Path) -> dict:
    local_dir = out_dir / "local"
    local_dir.mkdir(parents=True, exist_ok=True)

    proposal_spec = case["proposal_spec"]
    proposal_enriched = _enrich_partial_proposal_spec(proposal_spec)

    save_yaml(local_dir / "proposal-minimal.yaml", proposal_spec)
    save_yaml(local_dir / "proposal-enriched.yaml", proposal_enriched)
    save_yaml(local_dir / "diagram.yaml", case["diagram_spec"])
    save_yaml(local_dir / "business-case-minimal.yaml", case["business_case_spec"])
    save_yaml(local_dir / "bom.yaml", case["bom_spec"])

    deck_path = local_dir / f"{case['customer_id']}-deck.pptx"
    diagram_path = local_dir / f"{case['customer_id']}.drawio"
    bizcase_path = local_dir / f"{case['customer_id']}-bizcase.pptx"
    bom_path = local_dir / f"{case['customer_id']}-bom.xlsx"
    appca_path = local_dir / f"{case['customer_id']}-bom-appca.xlsx"
    pdf_path = local_dir / f"{case['customer_id']}.pdf"

    OCIDeckGenerator.from_spec(proposal_spec).save(str(deck_path))
    OCIDiagramGenerator.from_spec(case["diagram_spec"]).save(str(diagram_path))
    BusinessCaseDeckGenerator.from_spec(case["business_case_spec"]).save(str(bizcase_path))
    bom = OCIBomGenerator.from_spec(case["bom_spec"])
    bom.save(str(bom_path))
    bom.save(str(appca_path), appca=True)
    OCIPDFGenerator.from_spec(proposal_enriched, diagram_path=str(diagram_path)).save(str(pdf_path))

    analysis = {
        "deck": analyze_pptx(deck_path),
        "diagram": analyze_drawio(diagram_path),
        "business_case": analyze_pptx(bizcase_path),
        "bom": analyze_xlsx(bom_path),
        "bom_appca": analyze_xlsx(appca_path),
        "pdf": analyze_pdf(pdf_path),
    }
    return analysis


def main():
    parser = argparse.ArgumentParser(description="Generate random local evaluation cases for OCI Deal Accelerator.")
    parser.add_argument("--iterations", type=int, default=10, help="How many cases to generate.")
    parser.add_argument("--start-seed", type=int, default=1, help="Starting seed number.")
    parser.add_argument("--output-dir", default="tmp/evals", help="Directory for generated cases.")
    args = parser.parse_args()

    base_dir = PROJECT_ROOT / args.output_dir
    base_dir.mkdir(parents=True, exist_ok=True)

    manifest = []
    for seed in range(args.start_seed, args.start_seed + args.iterations):
        case = build_case(seed)
        case_dir = base_dir / f"iter-{seed:02d}"
        case_dir.mkdir(parents=True, exist_ok=True)
        (case_dir / "prompt.txt").write_text(case["prompt"], encoding="utf-8")
        (case_dir / "mcp-payloads.json").write_text(json.dumps(case["mcp_payloads"], indent=2), encoding="utf-8")
        local_analysis = render_local(case, case_dir)
        summary = {
            "seed": seed,
            "customer": case["customer"],
            "customer_id": case["customer_id"],
            "prompt": case["prompt"],
            "local_analysis": local_analysis,
        }
        (case_dir / "local-analysis.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
        manifest.append({
            "seed": seed,
            "customer": case["customer"],
            "customer_id": case["customer_id"],
            "prompt_path": str((case_dir / "prompt.txt").relative_to(PROJECT_ROOT)),
            "mcp_payloads_path": str((case_dir / "mcp-payloads.json").relative_to(PROJECT_ROOT)),
            "local_analysis_path": str((case_dir / "local-analysis.json").relative_to(PROJECT_ROOT)),
        })

    manifest_path = base_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(manifest_path.relative_to(PROJECT_ROOT))


if __name__ == "__main__":
    main()
